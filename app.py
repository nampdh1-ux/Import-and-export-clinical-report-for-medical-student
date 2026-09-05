import base64
import hashlib
import io
import json
import os
import random
import re
import smtplib
import tempfile
import time
from datetime import datetime
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

from fpdf import FPDF
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import google.generativeai as genai
import streamlit as st
from streamlit_local_storage import LocalStorage
from streamlit_pdf_viewer import pdf_viewer

# --- CẤU HÌNH TRANG ĐẦU TIÊN (Phải luôn nằm trên cùng) ---
st.set_page_config(page_title="Bệnh án Lâm sàng", layout="wide")

# ==============================================================================
# HÀM ĐIỀU PHỐI API KEY (CHỐNG RATE LIMIT)
# ==============================================================================
def get_feature_model(feature_key_name, model_name="gemini-3.1-flash-lite"):
    """Lấy model AI với API key chuyên biệt cho từng tác vụ."""
    api_key = st.secrets.get(feature_key_name) or st.secrets.get("GEMINI_API_KEY")
    if not api_key:
        return None
    genai.configure(api_key=api_key)
    return genai.GenerativeModel(model_name)

# ==============================================================================
# BẢO MẬT & XÁC THỰC DANH TÍNH (OTP + GMAIL + THIẾT BỊ)
# ==============================================================================
AUTH_STORAGE_KEY = "clinical_user_auth_token"
local_storage = LocalStorage()

def generate_auth_token(email):
    secret = st.secrets.get("AUTH_SECRET_KEY", "default_secret_medical_key_2026")
    raw_str = f"{email}_{secret}"
    return hashlib.sha256(raw_str.encode("utf-8")).hexdigest()

def send_otp_email(target_email, otp_code):
    sender_mail = st.secrets.get("SENDER_EMAIL")
    sender_pass = st.secrets.get("SENDER_APP_PASSWORD")
    if not (sender_mail and sender_pass):
        st.error("⚠️ Hệ thống chưa cấu hình SENDER_EMAIL hoặc SENDER_APP_PASSWORD trong Secrets!")
        return False
    try:
        msg = MIMEMultipart()
        msg['From'] = sender_mail
        msg['To'] = target_email
        msg['Subject'] = f"🔑 Mã xác thực truy cập Bệnh án Lâm sàng: {otp_code}"
        body = f"Xin chào,\n\nMã xác thực (OTP) dùng để đăng nhập vào Ứng dụng Bệnh án Lâm sàng của bạn là:\n\n👉  {otp_code}  👈\n\nMã có hiệu lực trong phiên đăng nhập này."
        msg.attach(MIMEText(body, 'plain'))
        server = smtplib.SMTP('smtp.gmail.com', 587, timeout=10)
        server.starttls()
        server.login(sender_mail, sender_pass)
        server.send_message(msg)
        server.quit()
        return True
    except Exception as e:
        st.error(f"Lỗi kết nối gửi email xác thực: {e}")
        return False

def send_login_notification(user_email):
    admin_mail = st.secrets.get("ADMIN_EMAIL")
    sender_mail = st.secrets.get("SENDER_EMAIL")
    sender_pass = st.secrets.get("SENDER_APP_PASSWORD")
    if not (admin_mail and sender_mail and sender_pass): return
    try:
        msg = MIMEMultipart()
        msg['From'] = sender_mail
        msg['To'] = admin_mail
        msg['Subject'] = f"🔔 [Bệnh Án Lâm Sàng] Người dùng mới đăng nhập: {user_email}"
        login_time = datetime.now().strftime('%d/%m/%Y %H:%M:%S')
        body = f"Hệ thống Bệnh Án Lâm Sàng ghi nhận lượt truy cập thành công:\n- Người dùng: {user_email}\n- Thời gian: {login_time}"
        msg.attach(MIMEText(body, 'plain'))
        server = smtplib.SMTP('smtp.gmail.com', 587, timeout=10)
        server.starttls()
        server.login(sender_mail, sender_pass)
        server.send_message(msg)
        server.quit()
    except Exception:
        pass

def check_password():
    admin_token_secret = str(st.secrets.get("ADMIN_BYPASS_TOKEN", "")).strip()
    url_admin_key = str(st.query_params.get("nam", "")).strip()
    is_admin_access = False
    if admin_token_secret and url_admin_key == admin_token_secret: is_admin_access = True
    elif "nam" in st.query_params and not admin_token_secret: is_admin_access = True

    if is_admin_access:
        st.session_state["password_correct"] = True
        st.session_state["is_admin"] = True
        if "logged_in_user" not in st.session_state: st.session_state["logged_in_user"] = "Admin"
        return True

    if st.session_state.get("password_correct"): return True

    try:
        saved_auth_raw = local_storage.getItem(AUTH_STORAGE_KEY)
        if saved_auth_raw:
            auth_data = json.loads(saved_auth_raw) if isinstance(saved_auth_raw, str) else saved_auth_raw
            if isinstance(auth_data, dict):
                saved_email = auth_data.get("email", "")
                saved_token = auth_data.get("token", "")
                if saved_email and saved_token == generate_auth_token(saved_email):
                    st.session_state["password_correct"] = True
                    st.session_state["logged_in_user"] = saved_email
                    if not st.session_state.get("sinh_vien"): st.session_state["sinh_vien"] = saved_email.split("@")[0]
                    st.rerun()
    except Exception: pass

    GMAIL_REGEX = r"^[a-zA-Z0-9](\.?[a-zA-Z0-9_-]){5,29}@gmail\.com$"
    with st.container():
        st.markdown("### 🔒 Ứng dụng Bệnh án Lâm sàng (Nội bộ)")
        st.caption("Vui lòng xác thực tài khoản Gmail chính chủ. Sau khi xác thực, thiết bị này sẽ được tự động ghi nhớ:")
        col_form, _ = st.columns([1.5, 1])
        with col_form:
            input_email = st.text_input("Địa chỉ Gmail của bạn:", placeholder="tenban@gmail.com", key="login_email_input").strip().lower()
            c_otp_btn, _ = st.columns([1, 1.5])
            with c_otp_btn: btn_send_otp = st.button("📩 Gửi mã xác thực OTP", use_container_width=True)
            
            if btn_send_otp:
                if not input_email or not re.match(GMAIL_REGEX, input_email):
                    st.error("❌ Vui lòng nhập địa chỉ Gmail hợp lệ (@gmail.com)!")
                else:
                    otp_random = str(random.randint(100000, 999999))
                    with st.spinner("Đang gửi mã xác thực về hộp thư của bạn..."):
                        if send_otp_email(input_email, otp_random):
                            st.session_state["generated_otp"] = otp_random
                            st.session_state["otp_target_email"] = input_email
                            st.success(f"✅ Đã gửi mã OTP đến {input_email}. Vui lòng mở hộp thư kiểm tra!")

            input_otp = st.text_input("Mã OTP (6 chữ số từ Gmail):", placeholder="VD: 123456", key="login_otp_input").strip()
            input_pass = st.text_input("Mã truy cập nội bộ:", type="password", key="login_pass_input")
            btn_login = st.button("Đăng nhập & Ghi nhớ máy này", type="primary", use_container_width=True)
            
            if btn_login:
                mat_khau_chuan = str(st.secrets.get("APP_PASSWORD", "123456")).strip()
                sent_otp = st.session_state.get("generated_otp")
                verified_email = st.session_state.get("otp_target_email")
                
                if not input_email or input_email != verified_email: st.error("❌ Email này chưa nhận mã OTP hoặc bị sửa đổi!")
                elif not input_otp or input_otp != sent_otp: st.error("❌ Mã OTP không chính xác!")
                elif input_pass != mat_khau_chuan: st.error("❌ Mã truy cập nội bộ không chính xác!")
                else:
                    with st.spinner("Đang lưu trạng thái và ghi nhớ thiết bị..."):
                        send_login_notification(input_email)
                        token = generate_auth_token(input_email)
                        auth_payload = json.dumps({"email": input_email, "token": token}, ensure_ascii=False)
                        local_storage.setItem(AUTH_STORAGE_KEY, auth_payload)
                        time.sleep(1.2)
                        st.session_state["password_correct"] = True
                        st.session_state["logged_in_user"] = input_email
                        if not st.session_state.get("sinh_vien"): st.session_state["sinh_vien"] = input_email.split("@")[0]
                        st.session_state.pop("generated_otp", None)
                        st.session_state.pop("otp_target_email", None)
                    st.toast("✅ Xác thực thành công!", icon="🎉")
                    st.rerun()
    return False

if not check_password(): st.stop()

# --- HÀM NÉN VÀ TỐI ƯU ẢNH PHIẾU XÉT NGHIỆM TRƯỚC KHI OCR ---
def optimize_lab_image(photo_file, max_dimension=1600, quality=85):
    try:
        photo_file.seek(0)
        img = Image.open(photo_file)
        if img.mode in ("RGBA", "P"): img = img.convert("RGB")
        width, height = img.size
        if max(width, height) > max_dimension:
            if width > height:
                new_width = max_dimension
                new_height = int(height * (max_dimension / width))
            else:
                new_height = max_dimension
                new_width = int(width * (max_dimension / height))
            img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
        buffer = io.BytesIO()
        img.save(buffer, format="JPEG", quality=quality, optimize=True)
        buffer.seek(0)
        return Image.open(buffer)
    except Exception:
        photo_file.seek(0)
        return Image.open(photo_file)

# ==============================================================================
# DANH MỤC TRƯỜNG DỮ LIỆU & NẠP BẢN NHÁP TỰ ĐỘNG
# ==============================================================================
STORAGE_KEY = "clinical_report_draft"

FIELDS_TO_SAVE = [
    "loai_benh_an",  # Phân loại bệnh án
    "ho_ten", "tuoi", "gioi_tinh", "dan_tok", "nghe_nghiep", "khoa_phong", "dia_chi", "ngay_vao_vien", "sinh_vien",
    "ly_do_vao_vien", "benh_su", 
    "bs_truoc_mo", "bs_trong_mo", "bs_sau_mo", 
    "ts_noi_khoa", "ts_ngoai_khoa", "ts_loi_song", "ts_gia_dinh",
    "kham_vao_vien", "kham_toan_than", "sh_mach", "sh_nhiet_do", "sh_ha", 
    "sh_nhip_tho", "sh_can_nang", "sh_chieu_cao", "sh_bmi", "sh_bmi_eval",
    "ngay_hau_phau", "kham_vet_mo", "kham_dan_luu", 
    "uu_tien_co_quan", "kham_tuan_hoan", "kham_ho_hap", "kham_tieu_hoa", 
    "kham_than_kinh", "kham_tiet_nieu", "kham_co_xuong_khop", "kham_co_quan_khac",
    "tom_tat", "chan_doan_so_bo", "chan_doan_phan_biet", "bien_luan",
    "cls_dx_xac_dinh", "cls_dx_dieu_tri", "cls_dx_khac",
    "chan_doan_xac_dinh", "bien_luan_xac_dinh",
    "dt_muc_tieu", "dt_cu_the", "dt_theo_doi",
    "tien_luong", "tu_van", "so_hang_cls"
]

def load_draft_to_session(loaded_ls):
    if "so_hang_cls" in loaded_ls: st.session_state["so_hang_cls"] = int(loaded_ls["so_hang_cls"])
    for k in FIELDS_TO_SAVE:
        if k in loaded_ls: st.session_state[k] = loaded_ls[k]
    try: st.session_state["tuoi"] = int(loaded_ls.get("tuoi", 45))
    except (ValueError, TypeError): st.session_state["tuoi"] = 45
    try: st.session_state["sh_can_nang"] = float(loaded_ls.get("sh_can_nang") or 0.0)
    except (ValueError, TypeError): st.session_state["sh_can_nang"] = 0.0
    try: st.session_state["sh_chieu_cao"] = float(loaded_ls.get("sh_chieu_cao") or 0.0)
    except (ValueError, TypeError): st.session_state["sh_chieu_cao"] = 0.0
    for i in range(st.session_state.get("so_hang_cls", 3)):
        if f"cls_kq_{i}" in loaded_ls: st.session_state[f"cls_kq_{i}"] = loaded_ls[f"cls_kq_{i}"]
        if f"cls_pg_{i}" in loaded_ls: st.session_state[f"cls_pg_{i}"] = loaded_ls[f"cls_pg_{i}"]

if "da_khoi_phuc_tu_dong" not in st.session_state:
    try:
        draft_raw = local_storage.getItem(STORAGE_KEY)
        if draft_raw:
            loaded_ls = json.loads(draft_raw) if isinstance(draft_raw, str) else draft_raw
            load_draft_to_session(loaded_ls)
    except Exception: pass
    st.session_state["da_khoi_phuc_tu_dong"] = True

# Khởi tạo giá trị mặc định
if "so_hang_cls" not in st.session_state: st.session_state["so_hang_cls"] = 3
for field in FIELDS_TO_SAVE:
    if field not in st.session_state:
        if field == "loai_benh_an": st.session_state[field] = "Nội khoa / Tiền phẫu"
        elif field == "bs_trong_mo":
            st.session_state[field] = (
                "- Hình thức mổ: Mổ phiên / Mổ cấp cứu\n"
                "- Phương pháp mổ: \n"
                "- Phương pháp gây mê: \n"
                "- Quá trình mổ: Không có biến chứng\n"
                "- Chẩn đoán sau mổ: "
            )
        elif field == "tuoi": st.session_state[field] = 45
        elif field == "gioi_tinh": st.session_state[field] = "Nam"
        elif field == "dan_tok": st.session_state[field] = "Kinh"
        elif field == "ngay_vao_vien": st.session_state[field] = datetime.now().strftime("%d/%m/%Y %H:%M")
        elif field in ["sh_can_nang", "sh_chieu_cao"]: st.session_state[field] = 0.0
        elif field == "uu_tien_co_quan": st.session_state[field] = "Không ưu tiên (Thứ tự mặc định)"
        else: st.session_state[field] = ""

for i in range(st.session_state["so_hang_cls"]):
    if f"cls_kq_{i}" not in st.session_state: st.session_state[f"cls_kq_{i}"] = ""
    if f"cls_pg_{i}" not in st.session_state: st.session_state[f"cls_pg_{i}"] = ""

# --- CSS TÙY BIẾN ---
st.markdown("""
<style>
    div[data-testid="stExpander"] { border: 1px solid #d4eaf0; border-radius: 6px; margin-bottom: 12px; background-color: #ffffff; }
    div[data-testid="stExpander"] > details > summary {
        background-color: #ebf7f9 !important; box-shadow: inset 4px 0 0 #c2185b, inset 8px 0 0 #0d47a1 !important;
        border-left: none !important; border-radius: 5px 5px 0 0; padding: 10px 14px 10px 18px !important;
        font-weight: 700 !important; color: #06445c !important; font-size: 1.05rem !important;
    }
    div[data-testid="stExpander"] > details > summary:hover { background-color: #ddf2f5 !important; color: #032b3b !important; }
    .sidebar-header-amboss { background-color: #ebf7f9; box-shadow: inset 4px 0 0 #c2185b, inset 8px 0 0 #0d47a1; padding: 8px 12px 8px 16px; border-radius: 4px; font-size: 1.05rem; font-weight: 700; color: #06445c; margin-bottom: 8px; }
    .sub-section-header { background-color: #f2fafb; box-shadow: inset 3px 0 0 #c2185b, inset 6px 0 0 #0d47a1; padding: 6px 12px 6px 14px; border-radius: 3px; margin-top: 10px; margin-bottom: 8px; font-size: 0.95rem; font-weight: 600; color: #0c4d63; }
    .highlight-dx { color: #b40000; font-weight: bold; font-size: 1.02rem; }
    .type-selector { padding: 15px; background-color: #fff9e6; border-left: 5px solid #ffc107; border-radius: 5px; margin-bottom: 20px;}
</style>
""", unsafe_allow_html=True)


# ==============================================================================
# HÀM HỖ TRỢ XUẤT FILE & AI CONTEXT
# ==============================================================================
def get_benh_su_text_for_ai():
    if st.session_state.get("loai_benh_an") == "Hậu phẫu":
        return f"- Trước mổ: {st.session_state.get('bs_truoc_mo')}\n- Trong mổ: {st.session_state.get('bs_trong_mo')}\n- Sau mổ: {st.session_state.get('bs_sau_mo')}"
    return st.session_state.get("benh_su")

def format_bullet_points(text):
    if not text or not str(text).strip(): return "Chưa ghi nhận thông tin."
    lines = str(text).strip().split("\n")
    formatted_lines = []
    for line in lines:
        cleaned = line.strip()
        if cleaned:
            if not cleaned.startswith("-") and not cleaned.startswith("*"): formatted_lines.append(f"- {cleaned}")
            else: formatted_lines.append(cleaned)
    return "\n".join(formatted_lines)

NORMAL_ORGAN_FINDINGS = {
    "kham_tuan_hoan": "- Lồng ngực cân đối, không ổ đập bất thường, không sẹo mổ cũ.\n- Mỏm tim đập ở khoang liên sườn V đường giữa đòn trái, diện đập 1-2 cm.\n- Dấu hiệu Hartzer (-), không có rung miêu.\n- Nhịp tim đều, tần số trùng nhịp mạch.\n- T1, T2 rõ, không nghe thấy tiếng tim bệnh lý (T3, T4, tiếng cọ màng ngoài tim).\n- Không có tiếng thổi bệnh lý ở các ổ van tim.\n- Mạch ngoại vi bắt rõ, đều hai bên.",
    "kham_ho_hap": "- Lồng ngực hai bên cân đối, di động đều theo nhịp thở, không co kéo cơ hô hấp phụ.\n- Khoang liên sườn không giãn rộng, không có tuần hoàn bàng hệ.\n- Rung thanh đều hai bên phế trường.\n- Gõ trong hai bên phổi.\n- Rì rào phế nang êm dịu hai phế trường.\n- Không nghe thấy rale ẩm, rale nổ, rale rít hay rale ngáy.",
    "kham_tieu_hoa": "- Bụng thon đều hai bên, di động theo nhịp thở, không chướng, không tuần hoàn bàng hệ, không sẹo mổ cũ.\n- Bụng mềm, không có điểm đau khu trú, không có phản ứng thành bụng hay cảm ứng phúc mạc.\n- Gan, lách không sờ thấy dưới bờ sườn, chiều cao gan trong giới hạn bình thường.\n- Các điểm đau ngoại khoa (Ruột thừa, Murphy, túi mật) âm tính.\n- Gõ trong toàn bụng, không có diện đục vùng thấp.\n- Tiếng nhu động ruột bình thường, không có tiếng thổi mạch máu bụng.",
    "kham_than_kinh": "- Bệnh nhân tỉnh táo, tiếp xúc tốt, Glasgow 15 điểm.\n- Không có dấu hiệu thần kinh khu trú.\n- Khám 12 đôi dây thần kinh sọ chưa phát hiện bệnh lý.\n- Trương lực cơ bình thường, cơ lực hai bên đều nhau (5/5).\n- Phản xạ gân xương tứ chi bình thường, đối xứng hai bên.\n- Dấu hiệu gáy mềm, Kernig (-), Brudzinski (-), Babinski (-) hai bên.\n- Cảm giác nông và sâu bình thường.",
    "kham_tiet_nieu": "- Hố thắt lưng hai bên cân đối, không sưng đỏ, không gồ cao.\n- Chạm thận (-), Bập bềnh thận (-).\n- Rung thận (-) hai bên.\n- Ấn các điểm niệu quản trên và giữa không đau.\n- Cầu bàng quang (-).",
    "kham_co_xuong_khop": "- Các khớp không sưng, nóng, đỏ, không biến dạng hay lệch trục.\n- Tầm vận động chủ động và thụ động các khớp trong giới hạn bình thường.\n- Không teo cơ, không cứng khớp buổi sáng.\n- Cột sống không gù vẹo, không có điểm đau chói dọc gai sống.",
    "kham_co_quan_khac": "- Răng - Hàm - Mặt, Tai - Mũi - Họng: Chưa phát hiện bất thường.\n- Nội tiết: Tuyến giáp không to, không có dấu hiệu suy hay cường giáp trên lâm sàng."
}

class BenhAnPDF(FPDF):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.add_font("Roboto", "", "Roboto-Regular.ttf")
        self.add_font("Roboto-Bold", "", "Roboto-Bold.ttf")
        self.loai_ba = "Nội khoa / Tiền phẫu"

    def header(self):
        if self.page_no() == 1:
            self.set_font("Roboto-Bold", "", 15)
            title = "BỆNH ÁN HẬU PHẪU" if self.loai_ba == "Hậu phẫu" else "BỆNH ÁN LÂM SÀNG"
            self.cell(0, 8, title, align="C", new_x="LMARGIN", new_y="NEXT")
            self.set_font("Roboto", "", 8)
            self.cell(0, 4, f"Thời gian làm bệnh án: {datetime.now().strftime('%d/%m/%Y %H:%M')}", align="C", new_x="LMARGIN", new_y="NEXT")
            self.ln(3)

    def footer(self):
        self.set_y(-12)
        self.set_font("Roboto", "", 8)
        self.cell(0, 10, f"Trang {self.page_no()}/{{nb}}", align="C")

    def add_section_header(self, title):
        self.set_font("Roboto-Bold", "", 11)
        self.set_fill_color(225, 235, 245)
        self.cell(0, 7, title, fill=True, new_x="LMARGIN", new_y="NEXT")
        self.ln(1)

    def add_subsection_header(self, title):
        self.set_font("Roboto-Bold", "", 10)
        self.cell(0, 6, title, new_x="LMARGIN", new_y="NEXT")

    def add_body_text(self, text):
        self.set_font("Roboto", "", 9.5)
        self.set_text_color(0, 0, 0)
        self.multi_cell(0, 5, str(text).strip() if str(text).strip() else "Chưa ghi nhận thông tin.")
        self.ln(2)

    def add_highlight_text(self, text):
        self.set_font("Roboto-Bold", "", 10.5)
        self.set_text_color(180, 0, 0)
        self.multi_cell(0, 5.5, str(text).strip() if str(text).strip() else "Chưa ghi nhận thông tin.")
        self.set_text_color(0, 0, 0)
        self.ln(2)

    def render_tom_tat_pdf(self, text):
        if not text or not str(text).strip():
            self.add_body_text("Chưa ghi nhận thông tin.")
            return
        lines = [line.strip() for line in str(text).strip().split("\n") if line.strip()]
        if not lines: return
        self.set_font("Roboto", "", 9.5)
        self.set_text_color(0, 0, 0)
        self.multi_cell(0, 5, lines[0])
        self.ln(1)
        orig_l_margin = self.l_margin
        self.set_left_margin(orig_l_margin + 8.0)
        for line in lines[1:]:
            bullet_line = line if (line.startswith("-") or line.startswith("*")) else f"- {line}"
            self.multi_cell(0, 5, bullet_line)
            self.ln(1)
        self.set_left_margin(orig_l_margin)
        self.ln(2)

    def render_table_cls(self, cls_rows):
        col_w = 95
        line_h = 5.0
        self.set_font("Roboto-Bold", "", 9.5)
        self.set_fill_color(230, 235, 245)
        if self.get_y() > 260: self.add_page()
        self.cell(col_w, 7, "KẾT QUẢ CẬN LÂM SÀNG", border=1, align="C", fill=True)
        self.cell(col_w, 7, "PHIÊN GIẢI / BIỆN GIẢI", border=1, align="C", fill=True, new_x="LMARGIN", new_y="NEXT")
        
        temp_files = []
        try:
            self.set_font("Roboto", "", 9)
            for kq, pg, img in cls_rows:
                txt_kq = format_bullet_points(kq) if kq else ""
                txt_pg = format_bullet_points(pg) if pg else "-"
                temp_img_path = None
                img_h = 0
                if img:
                    suffix = os.path.splitext(img.name)[1]
                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as t_img:
                        t_img.write(img.getbuffer())
                        temp_img_path = t_img.name
                        temp_files.append(temp_img_path)
                    img_h = 45
                nb_lines_left = len(self.multi_cell(col_w - 4, line_h, txt_kq, dry_run=True, output="LINES"))
                nb_lines_right = len(self.multi_cell(col_w - 4, line_h, txt_pg, dry_run=True, output="LINES"))
                
                h_left = nb_lines_left * line_h + 4 + (img_h + 4 if img else 0)
                h_right = nb_lines_right * line_h + 4
                row_h = max(h_left, h_right, 8)
                
                if self.get_y() + row_h > 275:
                    self.add_page()
                    self.set_font("Roboto-Bold", "", 9.5)
                    self.set_fill_color(230, 235, 245)
                    self.cell(col_w, 7, "KẾT QUẢ CẬN LÂM SÀNG", border=1, align="C", fill=True)
                    self.cell(col_w, 7, "PHIÊN GIẢI / BIỆN GIẢI", border=1, align="C", fill=True, new_x="LMARGIN", new_y="NEXT")
                    self.set_font("Roboto", "", 9)

                curr_x = self.get_x()
                curr_y = self.get_y()
                self.rect(curr_x, curr_y, col_w, row_h)
                self.rect(curr_x + col_w, curr_y, col_w, row_h)
                
                self.set_xy(curr_x + 2, curr_y + 2)
                if txt_kq: self.multi_cell(col_w - 4, line_h, txt_kq)
                if temp_img_path:
                    y_img = self.get_y() + 1
                    img_w = min(col_w - 8, 70)
                    x_img = curr_x + (col_w - img_w) / 2
                    try: self.image(temp_img_path, x=x_img, y=y_img, w=img_w, h=img_h)
                    except: pass
                
                self.set_xy(curr_x + col_w + 2, curr_y + 2)
                self.multi_cell(col_w - 4, line_h, txt_pg)
                self.set_xy(curr_x, curr_y + row_h)
            self.ln(3)
        finally:
            for p in temp_files:
                if os.path.exists(p):
                    try: os.remove(p)
                    except: pass

def export_pdf(data):
    pdf = BenhAnPDF()
    pdf.loai_ba = data.get("loai_benh_an", "Nội khoa / Tiền phẫu")
    pdf.alias_nb_pages()
    pdf.add_page()
    
    # I. HÀNH CHÍNH
    pdf.add_section_header("I. PHẦN HÀNH CHÍNH")
    hc_text = (
        f"- Họ và tên: {str(data['ho_ten']).upper()}    |    Tuổi: {data['tuoi']}    |    Giới tính: {data['gioi_tinh']}\n"
        f"- Dân tộc: {data['dan_tok']}    |    Nghề nghiệp: {data['nghe_nghiep']}\n"
        f"- Khoa phòng: {data['khoa_phong']}\n"
        f"- Địa chỉ: {data['dia_chi']}\n"
        f"- Ngày giờ vào viện: {data['ngay_vao_vien']}\n"
        f"- Sinh viên thực hiện: {data['sinh_vien']}"
    )
    pdf.add_body_text(hc_text)

    # II. LÝ DO VÀO VIỆN
    pdf.add_section_header("II. LÝ DO VÀO VIỆN")
    pdf.add_body_text(data['ly_do_vao_vien'])

    # III. BỆNH SỬ
    pdf.add_section_header("III. BỆNH SỬ")
    if pdf.loai_ba == "Hậu phẫu":
        pdf.add_subsection_header("1. Tình trạng trước mổ:")
        pdf.add_body_text(format_bullet_points(data.get('bs_truoc_mo', '')))
        pdf.add_subsection_header("2. Tình trạng trong mổ:")
        pdf.add_body_text(format_bullet_points(data.get('bs_trong_mo', '')))
        pdf.add_subsection_header("3. Quá trình sau mổ:")
        pdf.add_body_text(format_bullet_points(data.get('bs_sau_mo', '')))
    else:
        pdf.add_body_text(data.get('benh_su', ''))

    # IV. TIỀN SỬ
    pdf.add_section_header("IV. TIỀN SỬ")
    pdf.add_subsection_header("1. Tiền sử nội khoa:")
    pdf.add_body_text(format_bullet_points(data.get('ts_noi_khoa', '')))
    pdf.add_subsection_header("2. Tiền sử ngoại khoa và dị ứng:")
    pdf.add_body_text(format_bullet_points(data.get('ts_ngoai_khoa', '')))
    pdf.add_subsection_header("3. Lối sống và thói quen:")
    pdf.add_body_text(format_bullet_points(data.get('ts_loi_song', '')))
    pdf.add_subsection_header("4. Tiền sử gia đình:")
    pdf.add_body_text(format_bullet_points(data.get('ts_gia_dinh', '')))

    # V. THĂM KHÁM LÂM SÀNG
    pdf.add_section_header("V. THĂM KHÁM LÂM SÀNG")
    
    if pdf.loai_ba == "Hậu phẫu":
        pdf.add_subsection_header("1. Thăm khám hiện tại:")
        pdf.add_highlight_text(f"Hậu phẫu: {data.get('ngay_hau_phau', '...')}")
        pdf.add_body_text("a. Toàn thân:")
    else:
        pdf.add_subsection_header("1. Thăm khám lúc vào viện:")
        pdf.add_body_text(format_bullet_points(data.get('kham_vao_vien', '')))
        pdf.add_subsection_header("2. Thăm khám hiện tại:")
        pdf.add_body_text("a. Toàn thân:")

    pdf.add_body_text(format_bullet_points(data.get('kham_toan_than', '')))
    
    # BẢNG SINH HIỆU
    mach_val = data.get('sh_mach') or "--"
    nhiet_val = data.get('sh_nhiet_do') or "--"
    ha_val = data.get('sh_ha') or "--"
    nt_val = data.get('sh_nhip_tho') or "--"
    cn_val = data.get('sh_can_nang') if float(data.get('sh_can_nang', 0)) > 0 else "--"
    cc_val = data.get('sh_chieu_cao') if float(data.get('sh_chieu_cao', 0)) > 0 else "--"
    bmi_num = data.get('sh_bmi', '')
    bmi_txt = data.get('sh_bmi_eval', '')

    pdf.ln(1)
    pdf.set_draw_color(180, 180, 180)
    pdf.set_fill_color(245, 247, 250)
    col_w4 = (pdf.w - pdf.l_margin - pdf.r_margin) / 4.0
    pdf.set_font("Roboto-Bold", size=8.5)
    pdf.cell(col_w4, 5.5, f"Mạch: {mach_val} ck/phút", border=1, fill=True)
    pdf.cell(col_w4, 5.5, f"Nhiệt độ: {nhiet_val} °C", border=1, fill=True)
    pdf.cell(col_w4, 5.5, f"Huyết áp: {ha_val} mmHg", border=1, fill=True)
    pdf.cell(col_w4, 5.5, f"Nhịp thở: {nt_val} l/phút", border=1, fill=True, ln=True)
    
    col_w3_1 = col_w4
    col_w3_2 = col_w4
    col_w3_3 = col_w4 * 2.0
    bmi_display = f"BMI: {bmi_num} kg/m² ({bmi_txt})" if bmi_num else "BMI: --"
    pdf.cell(col_w3_1, 5.5, f"Chiều cao: {cc_val} cm", border=1)
    pdf.cell(col_w3_2, 5.5, f"Cân nặng: {cn_val} kg", border=1)
    pdf.cell(col_w3_3, 5.5, bmi_display, border=1, ln=True)
    pdf.ln(2)
    
    if pdf.loai_ba == "Hậu phẫu":
        pdf.add_body_text("b. Tình trạng vết mổ và dẫn lưu:")
        pdf.add_subsection_header("Vết mổ:")
        pdf.add_body_text(format_bullet_points(data.get('kham_vet_mo', '')))
        pdf.add_subsection_header("Ống dẫn lưu:")
        pdf.add_body_text(format_bullet_points(data.get('kham_dan_luu', '')))
        pdf.add_body_text("c. Các cơ quan:")
    else:
        pdf.add_body_text("b. Các cơ quan:")
        
    organ_list = [
        {"key": "kham_tuan_hoan", "name": "Tuần hoàn"},
        {"key": "kham_ho_hap", "name": "Hô hấp"},
        {"key": "kham_tieu_hoa", "name": "Tiêu hóa"},
        {"key": "kham_than_kinh", "name": "Thần kinh"},
        {"key": "kham_tiet_nieu", "name": "Thận - Tiết niệu"},
        {"key": "kham_co_xuong_khop", "name": "Cơ xương khớp"},
        {"key": "kham_co_quan_khac", "name": "Các cơ quan khác"},
    ]
    selected_organ = data.get("uu_tien_co_quan", "Không ưu tiên (Thứ tự mặc định)")
    if selected_organ != "Không ưu tiên (Thứ tự mặc định)":
        fav = next((it for it in organ_list if it["name"] == selected_organ), None)
        others = [it for it in organ_list if it["name"] != selected_organ]
        render_list = ([fav] + others) if fav else organ_list
    else:
        render_list = organ_list

    for org in render_list:
        content = format_bullet_points(data.get(org["key"], ""))
        title_text = f"{org['name']}:"
        pdf.add_subsection_header(title_text)
        text_w = pdf.get_string_width(title_text)
        x = pdf.get_x()
        y = pdf.get_y() - 1
        pdf.set_draw_color(50, 50, 50)
        pdf.set_line_width(0.3)
        pdf.line(x, y, x + text_w, y)
        pdf.add_body_text(content)

    # ĐỊNH NGHĨA SỐ LA MÃ ĐỘNG THEO LOẠI BỆNH ÁN
    if pdf.loai_ba == "Hậu phẫu":
        num_cdsb, num_cdpb, num_blsb, num_dxcls, num_cls, num_tt, num_cdxd, num_blxd = "VI", "VII", "VIII", "IX", "X", "XI", "XII", "XIII"
    else:
        num_tt, num_cdsb, num_cdpb, num_blsb, num_dxcls, num_cls, num_cdxd, num_blxd = "VI", "VII", "VIII", "IX", "X", "XI", "XII", "XIII"

    def pdf_tt():
        pdf.add_section_header(f"{num_tt}. TÓM TẮT BỆNH ÁN")
        pdf.render_tom_tat_pdf(data.get('tom_tat', ''))

    def pdf_cdsb():
        pdf.add_section_header(f"{num_cdsb}. CHẨN ĐOÁN SƠ BỘ")
        pdf.add_body_text(data.get('chan_doan_so_bo', ''))
        pdf.add_section_header(f"{num_cdpb}. CHẨN ĐOÁN PHÂN BIỆT")
        pdf.add_body_text(data.get('chan_doan_phan_biet', ''))
        bl = str(data.get('bien_luan', '')).strip()
        if bl:
            pdf.add_section_header(f"{num_blsb}. BIỆN LUẬN CHẨN ĐOÁN SƠ BỘ")
            pdf.add_body_text(bl)

    def pdf_cls():
        pdf.add_section_header(f"{num_dxcls}. ĐỀ XUẤT CẬN LÂM SÀNG")
        pdf.add_subsection_header("1. Phục vụ chẩn đoán xác định:")
        pdf.add_body_text(format_bullet_points(data.get('cls_dx_xac_dinh', '')))
        pdf.add_subsection_header("2. Phục vụ điều trị:")
        pdf.add_body_text(format_bullet_points(data.get('cls_dx_dieu_tri', '')))
        pdf.add_subsection_header("3. Cận lâm sàng khác:")
        pdf.add_body_text(format_bullet_points(data.get('cls_dx_khac', '')))

        pdf.add_section_header(f"{num_cls}. CẬN LÂM SÀNG ĐÃ CÓ")
        cls_rows = []
        so_hang = data.get("so_hang_cls", 3)
        for i in range(so_hang):
            kq = data.get(f"cls_kq_{i}", "").strip()
            pg = data.get(f"cls_pg_{i}", "").strip()
            img = data.get(f"cls_img_{i}", None)
            if kq or pg or img: cls_rows.append((kq, pg, img))
        if not cls_rows: pdf.add_body_text("Chưa ghi nhận kết quả cận lâm sàng.")
        else: pdf.render_table_cls(cls_rows)

    def pdf_cdxd():
        pdf.add_section_header(f"{num_cdxd}. CHẨN ĐOÁN XÁC ĐỊNH")
        pdf.add_highlight_text(format_bullet_points(data.get('chan_doan_xac_dinh', '')))
        noi_dung_bl_xd = str(data.get('bien_luan_xac_dinh', '')).strip()
        if noi_dung_bl_xd:
            pdf.add_section_header(f"{num_blxd}. BIỆN LUẬN CHẨN ĐOÁN XÁC ĐỊNH")
            pdf.add_body_text(format_bullet_points(noi_dung_bl_xd))

    # TRIỂN KHAI TRẬT TỰ ĐỘNG
    if pdf.loai_ba == "Hậu phẫu":
        pdf_cdsb()
        pdf_cls()
        pdf_tt()
        pdf_cdxd()
    else:
        pdf_tt()
        pdf_cdsb()
        pdf_cls()
        pdf_cdxd()

    pdf.add_section_header("XIV. ĐIỀU TRỊ")
    pdf.add_subsection_header("1. Mục tiêu điều trị:")
    pdf.add_body_text(format_bullet_points(data.get('dt_muc_tieu', '')))
    pdf.add_subsection_header("2. Điều trị cụ thể:")
    pdf.add_body_text(format_bullet_points(data.get('dt_cu_the', '')))
    pdf.add_subsection_header("3. Theo dõi sau điều trị:")
    pdf.add_body_text(format_bullet_points(data.get('dt_theo_doi', '')))

    noi_dung_tien_luong = str(data.get("tien_luong", "")).strip()
    if noi_dung_tien_luong:
        pdf.add_section_header("XV. TIÊN LƯỢNG")
        pdf.add_body_text(format_bullet_points(noi_dung_tien_luong))

    noi_dung_tu_van = str(data.get("tu_van", "")).strip()
    if noi_dung_tu_van:
        ten_de_muc_tu_van = "XVI. TƯ VẤN" if noi_dung_tien_luong else "XV. TƯ VẤN"
        pdf.add_section_header(ten_de_muc_tu_van)
        pdf.add_body_text(format_bullet_points(noi_dung_tu_van))

    return bytes(pdf.output())

def export_pptx(data):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    loai_ba = data.get("loai_benh_an", "Nội khoa / Tiền phẫu")

    COLOR_PRIMARY = RGBColor(13, 71, 161)
    COLOR_ACCENT = RGBColor(194, 24, 91)
    COLOR_TEXT = RGBColor(30, 41, 59)
    COLOR_RED = RGBColor(180, 0, 0)

    def add_slide_with_header(title_text):
        slide = prs.slides.add_slide(blank_layout)
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Calibri"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        line_rose = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(0.6), Inches(0.06))
        line_rose.fill.solid()
        line_rose.fill.fore_color.rgb = COLOR_ACCENT
        line_rose.line.fill.background()
        
        line_blue = slide.shapes.add_shape(1, Inches(1.45), Inches(1.3), Inches(11.083), Inches(0.06))
        line_blue.fill.solid()
        line_blue.fill.fore_color.rgb = COLOR_PRIMARY
        line_blue.line.fill.background()
        return slide

    def add_content_with_overflow(title_text, items_list, is_red=False):
        if not items_list: return
        MAX_LINES_PER_SLIDE = 7
        total_chunks = [items_list[i:i + MAX_LINES_PER_SLIDE] for i in range(0, len(items_list), MAX_LINES_PER_SLIDE)]
        for idx, chunk in enumerate(total_chunks):
            current_title = title_text if idx == 0 else f"{title_text} (tiếp theo)"
            slide = add_slide_with_header(current_title)
            box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
            tf = box.text_frame
            tf.word_wrap = True
            
            for line_idx, item in enumerate(chunk):
                if isinstance(item, tuple): text, is_sub_header = item
                else: text, is_sub_header = item, False
                p = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
                p.text = text
                p.font.name = "Calibri"
                if is_sub_header:
                    p.font.size = Pt(19)
                    p.font.bold = True
                    p.font.underline = True
                    p.font.color.rgb = COLOR_PRIMARY
                    p.space_after = Pt(6)
                else:
                    p.font.size = Pt(16.5)
                    p.font.bold = is_red
                    p.font.color.rgb = COLOR_RED if is_red else COLOR_TEXT
                    p.space_after = Pt(10)

    # Bìa
    title_slide = prs.slides.add_slide(blank_layout)
    t_box = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf_t = t_box.text_frame
    p1 = tf_t.paragraphs[0]
    p1.text = "BỆNH ÁN HẬU PHẪU" if loai_ba == "Hậu phẫu" else "BỆNH ÁN LÂM SÀNG"
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(16)
    
    p2 = tf_t.add_paragraph()
    p2.text = f"Bệnh nhân: {str(data['ho_ten']).upper()} | {data['tuoi']} tuổi | Giới tính: {data['gioi_tinh']}"
    p2.font.size = Pt(20)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(8)

    p3 = tf_t.add_paragraph()
    p3.text = f"Khoa phòng: {data['khoa_phong']} | Người thực hiện: {data['sinh_vien']}"
    p3.font.size = Pt(16)
    p3.alignment = PP_ALIGN.CENTER

    # I. HÀNH CHÍNH
    hc_items = [
        (f"Họ và tên: {str(data['ho_ten']).upper()}", False),
        (f"Tuổi: {data['tuoi']}   |   Giới tính: {data['gioi_tinh']}   |   Dân tộc: {data['dan_tok']}", False),
        (f"Nghề nghiệp: {data['nghe_nghiep']}  |  Khoa / Phòng: {data['khoa_phong']}", False),
        (f"Địa chỉ: {data['dia_chi']}", False),
        (f"Ngày giờ vào viện: {data['ngay_vao_vien']}", False)
    ]
    add_content_with_overflow("I. PHẦN HÀNH CHÍNH", hc_items)

    # II & III
    bs_items = [("1. Lý do vào viện:", True), (f"- {data.get('ly_do_vao_vien', '')}", False)]
    if loai_ba == "Hậu phẫu":
        for title, key in [("2. Tình trạng trước mổ:", "bs_truoc_mo"), ("3. Tình trạng trong mổ:", "bs_trong_mo"), ("4. Quá trình sau mổ:", "bs_sau_mo")]:
            items = [l.strip() for l in str(data.get(key, "")).split("\n") if l.strip()]
            if items:
                bs_items.append((title, True))
                bs_items.extend([(f"- {l}" if not l.startswith("-") else l, False) for l in items])
    else:
        bs_text_lines = [l.strip() for l in str(data.get('benh_su', '')).split("\n") if l.strip()]
        if bs_text_lines:
            bs_items.append(("2. Bệnh sử:", True))
            bs_items.extend([(f"- {l}" if not l.startswith("-") else l, False) for l in bs_text_lines])
    add_content_with_overflow("II VÀ III. LÝ DO VÀO VIỆN VÀ BỆNH SỬ", bs_items)

    # IV
    ts_items = []
    for label, key in [("1. Tiền sử nội khoa:", "ts_noi_khoa"), ("2. Tiền sử ngoại khoa & dị ứng:", "ts_ngoai_khoa"), ("3. Lối sống & thói quen:", "ts_loi_song"), ("4. Tiền sử gia đình:", "ts_gia_dinh")]:
        lines = [l.strip() for l in str(data.get(key, "")).split("\n") if l.strip()]
        if lines:
            ts_items.append((label, True))
            ts_items.extend([(f"- {l}" if not l.startswith("-") else l, False) for l in lines])
    if ts_items: add_content_with_overflow("IV. TIỀN SỬ", ts_items)

    # V
    tk_items = []
    if loai_ba != "Hậu phẫu" and data.get("kham_vao_vien"): 
        tk_items.append(("1. Khám lúc vào viện:", True))
        tk_items.extend([(f"- {l.strip()}", False) for l in str(data['kham_vao_vien']).split("\n") if l.strip()])
        
    if data.get("kham_toan_than"):
        if loai_ba == "Hậu phẫu":
            tk_items.append(("1. Thăm khám hiện tại - Toàn thân:", True))
            tk_items.append((f"Hậu phẫu ngày thứ: {data.get('ngay_hau_phau', '...')}", False))
        else:
            tk_items.append(("2. Thăm khám hiện tại - Toàn thân:", True))
        tk_items.extend([(f"- {l.strip()}", False) for l in str(data['kham_toan_than']).split("\n") if l.strip()])
        
    if loai_ba == "Hậu phẫu" and (data.get("kham_vet_mo") or data.get("kham_dan_luu")):
        tk_items.append(("2. Vết mổ & Dẫn lưu:", True))
        if data.get("kham_vet_mo"): tk_items.append((f"- Vết mổ: {data['kham_vet_mo']}", False))
        if data.get("kham_dan_luu"): tk_items.append((f"- Ống dẫn lưu: {data['kham_dan_luu']}", False))
        tk_items.append(("3. Khám các cơ quan:", True))
    else:
        tk_items.append(("2. Khám các cơ quan:" if loai_ba == "Hậu phẫu" else "3. Khám các cơ quan:", True))
        
    cq_list = [("Tuần hoàn", "kham_tuan_hoan"), ("Hô hấp", "kham_ho_hap"), ("Tiêu hóa", "kham_tieu_hoa"), ("Thần kinh", "kham_than_kinh"), ("Thận - Tiết niệu", "kham_tiet_nieu"), ("Cơ xương khớp", "kham_co_xuong_khop")]
    for name, key in cq_list:
        val = str(data.get(key, "")).strip()
        if val: tk_items.append((f"- {name}: {val}", False))
    if tk_items: add_content_with_overflow("V. THĂM KHÁM LÂM SÀNG", tk_items)

    # ĐỊNH NGHĨA SỐ LA MÃ ĐỘNG PPTX
    if loai_ba == "Hậu phẫu":
        num_cdsb, num_cdpb, num_blsb, num_dxcls, num_cls, num_tt, num_cdxd, num_blxd = "VI", "VII", "VIII", "IX", "X", "XI", "XII", "XIII"
    else:
        num_tt, num_cdsb, num_cdpb, num_blsb, num_dxcls, num_cls, num_cdxd, num_blxd = "VI", "VII", "VIII", "IX", "X", "XI", "XII", "XIII"

    def pptx_tt():
        tt_items = []
        lines_tt = [l.strip() for l in str(data.get("tom_tat", "")).split("\n") if l.strip()]
        if lines_tt:
            tt_items.append(("Tóm tắt diễn biến ca bệnh:", True))
            tt_items.extend([(f"- {l}" if not l.startswith("-") else l, False) for l in lines_tt])
            add_content_with_overflow(f"{num_tt}. TÓM TẮT BỆNH ÁN", tt_items)

    def pptx_cdsb():
        cd_items = []
        if str(data.get("chan_doan_so_bo", "")).strip():
            cd_items.append(("Chẩn đoán sơ bộ:", True))
            cd_items.extend([(f"- {l.strip()}", False) for l in str(data['chan_doan_so_bo']).split("\n") if l.strip()])
        if str(data.get("chan_doan_phan_biet", "")).strip():
            cd_items.append(("Chẩn đoán phân biệt:", True))
            cd_items.extend([(f"- {l.strip()}", False) for l in str(data['chan_doan_phan_biet']).split("\n") if l.strip()])
        if cd_items: add_content_with_overflow(f"{num_cdsb} VÀ {num_cdpb}. CHẨN ĐOÁN SƠ BỘ VÀ PHÂN BIỆT", cd_items)

        lines_bl = [l.strip() for l in str(data.get("bien_luan", "")).split("\n") if l.strip()]
        if lines_bl:
            add_content_with_overflow(f"{num_blsb}. BIỆN LUẬN CHẨN ĐOÁN SƠ BỘ", [("Biện luận lâm sàng:", True)] + [(f"- {l}" if not l.startswith("-") else l, False) for l in lines_bl])

    def pptx_cls():
        cls_items = []
        for label, key in [("1. Phục vụ chẩn đoán xác định:", "cls_dx_xac_dinh"), ("2. Phục vụ điều trị:", "cls_dx_dieu_tri"), ("3. Cận lâm sàng khác:", "cls_dx_khac")]:
            lines = [l.strip() for l in str(data.get(key, "")).split("\n") if l.strip()]
            if lines:
                cls_items.append((label, True))
                cls_items.extend([(f"- {l}" if not l.startswith("-") else l, False) for l in lines])
        if cls_items: add_content_with_overflow(f"{num_dxcls}. ĐỀ XUẤT CẬN LÂM SÀNG", cls_items)

        cls_rows = []
        so_hang = data.get("so_hang_cls", 3)
        for i in range(so_hang):
            kq = data.get(f"cls_kq_{i}", "").strip()
            pg = data.get(f"cls_pg_{i}", "").strip()
            img = data.get(f"cls_img_{i}", None)
            if kq or pg or img: cls_rows.append((kq, pg, img))

        if cls_rows:
            rows_text_only = [item for item in cls_rows if not item[2]]
            rows_with_img = [item for item in cls_rows if item[2]]
            if rows_text_only:
                table_chunks = [rows_text_only[i:i + 3] for i in range(0, len(rows_text_only), 3)]
                for c_idx, chunk in enumerate(table_chunks):
                    slide = add_slide_with_header(f"{num_cls}. CẬN LÂM SÀNG ĐÃ CÓ" if c_idx == 0 else f"{num_cls}. CẬN LÂM SÀNG ĐÃ CÓ (tiếp theo)")
                    table_shape = slide.shapes.add_table(len(chunk) + 1, 2, Inches(0.8), Inches(1.6), Inches(11.733), Inches(1.0 + len(chunk) * 1.3))
                    table = table_shape.table
                    table.columns[0].width = Inches(5.866)
                    table.columns[1].width = Inches(5.866)
                    table.cell(0, 0).text = "KẾT QUẢ CẬN LÂM SÀNG"
                    table.cell(0, 1).text = "PHIÊN GIẢI / BIỆN GIẢI"
                    for col_i in range(2):
                        cell_p = table.cell(0, col_i).text_frame.paragraphs[0]
                        cell_p.font.bold = True
                        cell_p.font.size = Pt(14)
                        cell_p.font.color.rgb = COLOR_PRIMARY
                    for r_i, (kq, pg, _) in enumerate(chunk):
                        table.cell(r_i + 1, 0).text = kq if kq else "-"
                        table.cell(r_i + 1, 1).text = pg if pg else "-"
                        for col_i in range(2):
                            cell_p = table.cell(r_i + 1, col_i).text_frame.paragraphs[0]
                            cell_p.font.size = Pt(13)
            temp_img_files = []
            try:
                for kq, pg, img in rows_with_img:
                    slide = add_slide_with_header(f"{num_cls}. CẬN LÂM SÀNG ĐÃ CÓ (HÌNH ẢNH)")
                    suffix = os.path.splitext(img.name)[1]
                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as t_img:
                        t_img.write(img.getbuffer())
                        temp_path = t_img.name
                        temp_img_files.append(temp_path)
                    try: slide.shapes.add_picture(temp_path, Inches(0.8), Inches(1.6), width=Inches(5.6))
                    except: pass
                    tf_r = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.0)).text_frame
                    tf_r.word_wrap = True
                    p_kq_title = tf_r.paragraphs[0]
                    p_kq_title.text, p_kq_title.font.bold, p_kq_title.font.underline, p_kq_title.font.size, p_kq_title.font.color.rgb = "Kết quả ghi nhận:", True, True, Pt(16), COLOR_PRIMARY
                    p_kq = tf_r.add_paragraph()
                    p_kq.text, p_kq.font.size, p_kq.space_after = kq if kq else "Hình ảnh xét nghiệm đính kèm", Pt(14.5), Pt(16)
                    p_pg_title = tf_r.add_paragraph()
                    p_pg_title.text, p_pg_title.font.bold, p_pg_title.font.underline, p_pg_title.font.size, p_pg_title.font.color.rgb = "Biện giải / Phiên giải:", True, True, Pt(16), COLOR_PRIMARY
                    p_pg = tf_r.add_paragraph()
                    p_pg.text, p_pg.font.size = pg if pg else "-", Pt(14.5)
            finally:
                for p in temp_img_files:
                    try: os.remove(p)
                    except: pass

    def pptx_cdxd():
        cdxd_lines = [l.strip() for l in str(data.get("chan_doan_xac_dinh", "")).split("\n") if l.strip()]
        if cdxd_lines: add_content_with_overflow(f"{num_cdxd}. CHẨN ĐOÁN XÁC ĐỊNH", [("Chẩn đoán xác định:", True)] + [(f"- {l}" if not l.startswith("-") else l, False) for l in cdxd_lines], is_red=True)

        blxd_lines = [l.strip() for l in str(data.get("bien_luan_xac_dinh", "")).split("\n") if l.strip()]
        if blxd_lines: add_content_with_overflow(f"{num_blxd}. BIỆN LUẬN CHẨN ĐOÁN XÁC ĐỊNH", [("Biện luận chẩn đoán xác định:", True)] + [(f"- {l}" if not l.startswith("-") else l, False) for l in blxd_lines])

    # THỰC THI PPTX TRẬT TỰ ĐỘNG
    if loai_ba == "Hậu phẫu":
        pptx_cdsb()
        pptx_cls()
        pptx_tt()
        pptx_cdxd()
    else:
        pptx_tt()
        pptx_cdsb()
        pptx_cls()
        pptx_cdxd()

    dt_items = []
    for label, key in [("1. Mục tiêu điều trị:", "dt_muc_tieu"), ("2. Điều trị cụ thể:", "dt_cu_the"), ("3. Theo dõi sau điều trị:", "dt_theo_doi")]:
        lines = [l.strip() for l in str(data.get(key, "")).split("\n") if l.strip()]
        if lines:
            dt_items.append((label, True))
            dt_items.extend([(f"- {l}" if not l.startswith("-") else l, False) for l in lines])
    if dt_items: add_content_with_overflow("XIV. ĐIỀU TRỊ", dt_items)

    tl_lines = [l.strip() for l in str(data.get("tien_luong", "")).split("\n") if l.strip()]
    if tl_lines: add_content_with_overflow("XV. TIÊN LƯỢNG", [("Đánh giá tiên lượng bệnh nhân:", True)] + [(f"- {l}" if not l.startswith("-") else l, False) for l in tl_lines])

    tv_lines = [l.strip() for l in str(data.get("tu_van", "")).split("\n") if l.strip()]
    if tv_lines: add_content_with_overflow("XVI. TƯ VẤN" if tl_lines else "XV. TƯ VẤN", [("Hướng dẫn và tư vấn cho người bệnh:", True)] + [(f"- {l}" if not l.startswith("-") else l, False) for l in tv_lines])

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io.getvalue()

# ==============================================================================
# SIDEBAR: QUẢN LÝ BẢN NHÁP & ĐĂNG XUẤT
# ==============================================================================
with st.sidebar:
    st.markdown("<div class='sidebar-header-amboss'>Quản lý bản nháp</div>", unsafe_allow_html=True)
    st.caption("🟢 **Tự động lưu:** Dữ liệu được ghi nhớ tự động vào trình duyệt mỗi khi nhập liệu.")
    
    if st.button("🔄 Nạp lại bản nháp từ trình duyệt", type="primary", use_container_width=True):
        saved_raw = local_storage.getItem(STORAGE_KEY)
        if saved_raw:
            try:
                loaded_ls = json.loads(saved_raw) if isinstance(saved_raw, str) else saved_raw
                load_draft_to_session(loaded_ls)
                st.toast("Đã khôi phục bệnh án thành công!", icon="✅")
                st.rerun()
            except Exception as e: st.error(f"Lỗi khi đọc bản nháp: {e}")
        else: st.warning("Không tìm thấy dữ liệu nháp nào.")

    if st.button("🗑️ Xóa bản nháp (Làm bệnh án mới)", use_container_width=True):
        local_storage.deleteItem(STORAGE_KEY)
        for k in FIELDS_TO_SAVE:
            if k == "tuoi": st.session_state[k] = 45
            elif k in ["sh_can_nang", "sh_chieu_cao"]: st.session_state[k] = 0.0
            elif k == "gioi_tinh": st.session_state[k] = "Nam"
            elif k == "dan_tok": st.session_state[k] = "Kinh"
            elif k == "loai_benh_an": st.session_state[k] = "Nội khoa / Tiền phẫu"
            elif k == "bs_trong_mo":
                st.session_state[k] = (
                    "- Hình thức mổ: Mổ phiên / Mổ cấp cứu\n"
                    "- Phương pháp mổ: \n"
                    "- Phương pháp gây mê: \n"
                    "- Quá trình mổ: Không có biến chứng\n"
                    "- Chẩn đoán sau mổ: "
                )
            elif k == "uu_tien_co_quan": st.session_state[k] = "Không ưu tiên (Thứ tự mặc định)"
            elif k == "ngay_vao_vien": st.session_state[k] = datetime.now().strftime("%d/%m/%Y %H:%M")
            else: st.session_state[k] = ""
        st.session_state["so_hang_cls"] = 3
        for i in range(10): st.session_state[f"cls_kq_{i}"], st.session_state[f"cls_pg_{i}"] = "", ""
        st.session_state["last_saved_snapshot"] = ""
        st.toast("Đã xóa sạch bản nháp và làm mới form!", icon="🗑️")
        st.rerun()
        st.session_state["so_hang_cls"] = 3
        for i in range(10): st.session_state[f"cls_kq_{i}"], st.session_state[f"cls_pg_{i}"] = "", ""
        st.session_state["last_saved_snapshot"] = ""
        st.toast("Đã xóa sạch bản nháp và làm mới form!", icon="🗑️")
        st.rerun()

    st.markdown("---")
    st.caption("Hoặc lưu trữ dạng tập tin JSON tải về máy:")
    current_data = {k: st.session_state.get(k, "") for k in FIELDS_TO_SAVE}
    for i in range(current_data.get("so_hang_cls", 3)):
        current_data[f"cls_kq_{i}"] = st.session_state.get(f"cls_kq_{i}", "")
        current_data[f"cls_pg_{i}"] = st.session_state.get(f"cls_pg_{i}", "")
        
    json_string = json.dumps(current_data, ensure_ascii=False, indent=2)
    ten_benh_nhan = str(st.session_state.get("ho_ten", "chua_dat_ten")).strip().replace(" ", "_")
    if not ten_benh_nhan: ten_benh_nhan = "chua_dat_ten"
        
    st.download_button("📥 Lưu bản nháp về máy (.json)", data=json_string, file_name=f"Ban_nhap_{ten_benh_nhan}_{datetime.now().strftime('%Y%m%d_%H%M')}.json", mime="application/json", use_container_width=True)
    st.markdown("---")
    st.markdown("**Khôi phục dữ liệu từ bản nháp:**")
    file_nhap = st.file_uploader("Chọn tập tin .json đã lưu:", type=["json"], key="uploader_nhap_json")
    if file_nhap is not None:
        if st.button("🔄 Nhấn vào đây để nạp dữ liệu", type="primary", use_container_width=True):
            try:
                loaded_data = json.load(file_nhap)
                load_draft_to_session(loaded_data)
                st.success("Đã nạp bản nháp thành công!")
                st.rerun()
            except Exception as e: st.error(f"Không thể đọc file: {e}")

    if st.session_state.get("password_correct") and not st.session_state.get("is_admin"):
        st.markdown("---")
        if st.button("🚪 Đăng xuất khỏi thiết bị này"):
            local_storage.deleteItem(AUTH_STORAGE_KEY)
            st.session_state["password_correct"] = False
            st.session_state.pop("logged_in_user", None)
            st.rerun()

# ==============================================================================
# GIAO DIỆN CHÍNH (3 TABS)
# ==============================================================================
st.title("Bệnh Án Lâm Sàng")
st.caption("Cấu trúc bệnh án trình bày ca bệnh và thi lâm sàng (Hỗ trợ Nội khoa, Ngoại khoa, Hậu phẫu).")

loai_benh_an = st.radio("📌 **LỰA CHỌN MẪU BỆNH ÁN:**", ["Nội khoa / Tiền phẫu", "Hậu phẫu"], horizontal=True, key="loai_benh_an")

# Khai báo Dictionary lưu trữ ảnh toàn cục
uploaded_imgs = {}

# CÁC HÀM UI RỜI RẠC DÙNG CHUNG (Để hoán đổi vị trí)
# --- HÀM TỰ ĐỘNG TẠO CÂU DẪN TÓM TẮT BỆNH ÁN NỘI KHOA (KHÔNG DÙNG AI) ---
def generate_intro_tom_tat_noi_khoa():
    gioi_tinh = st.session_state.get("gioi_tinh", "Nam")
    tuoi = st.session_state.get("tuoi", "")
    tuoi_str = f"{tuoi} tuổi" if tuoi else ""
    
    ts_list = []
    for k in ["ts_noi_khoa", "ts_ngoai_khoa"]:
        val = str(st.session_state.get(k, "")).strip()
        if val:
            lines = [l.strip().lstrip("-*• ") for l in val.split("\n") if l.strip()]
            if lines:
                ts_list.append(", ".join(lines))
    tien_su_str = "; ".join(ts_list) if ts_list else "chưa ghi nhận bất thường"

    ly_do = str(st.session_state.get("ly_do_vao_vien", "")).strip() or "..."

    ngay_vv_raw = str(st.session_state.get("ngay_vao_vien", "")).strip()
    so_ngay_nam_vien = 0
    if ngay_vv_raw:
        try:
            date_match = re.search(r"(\d{1,2}/\d{1,2}/\d{4})", ngay_vv_raw)
            if date_match:
                d_vv = datetime.strptime(date_match.group(1), "%d/%m/%Y").date()
                d_hientai = datetime.now().date()
                so_ngay_nam_vien = max(0, (d_hientai - d_vv).days)
        except Exception:
            so_ngay_nam_vien = 0

    benh_su = str(st.session_state.get("benh_su", "")).lower()
    so_ngay_truoc_vv = 0
    match_ngay = re.search(r"cách (?:vào viện|nhập viện)\s*(?:khoảng|tầm)?\s*(\d+)\s*ngày", benh_su)
    match_gio = re.search(r"cách (?:vào viện|nhập viện)\s*(?:khoảng|tầm)?\s*(\d+)\s*giờ", benh_su)
    match_tuan = re.search(r"cách (?:vào viện|nhập viện)\s*(?:khoảng|tầm)?\s*(\d+)\s*tuần", benh_su)

    if match_ngay:
        so_ngay_truoc_vv = int(match_ngay.group(1))
    elif match_gio:
        so_ngay_truoc_vv = max(1, round(int(match_gio.group(1)) / 24))
    elif match_tuan:
        so_ngay_truoc_vv = int(match_tuan.group(1)) * 7

    tong_ngay = so_ngay_truoc_vv + so_ngay_nam_vien
    if tong_ngay > 0:
        dien_bien_str = f"bệnh diễn biến {tong_ngay} ngày nay"
    else:
        dien_bien_str = "bệnh diễn biến cấp tính"

    return (
        f"Bệnh nhân {gioi_tinh.lower()} {tuoi_str}, tiền sử {tien_su_str} vào viện vì {ly_do}, "
        f"{dien_bien_str}. Qua thăm khám và hỏi bệnh phát hiện các hội chứng và triệu chứng sau:"
    )

# --- HÀM TỰ ĐỘNG TẠO CÂU DẪN TÓM TẮT BỆNH ÁN HẬU PHẪU (KHÔNG DÙNG AI) ---
def generate_intro_tom_tat_hau_phau():
    gioi_tinh = st.session_state.get("gioi_tinh", "Nam")
    tuoi = st.session_state.get("tuoi", "")
    tuoi_str = f"{tuoi} tuổi" if tuoi else "..."
    ly_do = str(st.session_state.get("ly_do_vao_vien", "")).strip() or "..."

    bs_truoc_mo = str(st.session_state.get("bs_truoc_mo", "")).strip()
    bs_trong_mo = str(st.session_state.get("bs_trong_mo", "")).strip()
    cd_so_bo = str(st.session_state.get("chan_doan_so_bo", "")).strip()

    # 1. Trích xuất Chẩn đoán trước mổ
    cd_truoc_mo = "..."
    m_cdtm = re.search(r"(?:chẩn đoán trước mổ|cđ trước mổ)(?:\s*là)?[:\s\-]+([^.\n;]+)", bs_truoc_mo, re.IGNORECASE)
    if m_cdtm:
        cd_truoc_mo = m_cdtm.group(1).strip()
    elif bs_truoc_mo:
        cd_truoc_mo = bs_truoc_mo.split("\n")[-1].strip().lstrip("-*• ")

    # 2. Trích xuất Phương pháp mổ & Cấp cứu hay Chương trình/Phiên
    pp_mo = "..."
    loai_mo = "cấp cứu/phiên"
    if bs_trong_mo:
        m_pp = re.search(r"(?:phương pháp phẫu thuật|phương pháp mổ|phẫu thuật|pt)(?:\s*là)?[:\s\-]+([^.\n;]+)", bs_trong_mo, re.IGNORECASE)
        if m_pp:
            pp_mo = m_pp.group(1).strip()
        else:
            first_line = bs_trong_mo.split("\n")[0].strip().lstrip("-*• ")
            pp_mo = first_line if len(first_line) < 60 else "phẫu thuật"

        bs_tm_lower = bs_trong_mo.lower()
        if "cấp cứu" in bs_tm_lower:
            loai_mo = "cấp cứu"
        elif "chương trình" in bs_tm_lower or "mổ phiên" in bs_tm_lower:
            loai_mo = "chương trình"

    # 3. Trích xuất Chẩn đoán sau mổ
    cd_sau_mo = "..."
    m_cdsm = re.search(r"(?:chẩn đoán sau mổ|cđ sau mổ)(?:\s*là)?[:\s\-]+([^.\n;]+)", bs_trong_mo, re.IGNORECASE)
    if m_cdsm:
        cd_sau_mo = m_cdsm.group(1).strip()
    elif cd_so_bo:
        # Nếu đã ghi ở CĐ sơ bộ: "Hậu phẫu ngày 3 - Viêm phúc mạc ruột thừa..." -> lấy phần bệnh học
        cd_clean = re.sub(r"^hậu phẫu ngày[^-\:]*[-\:]\s*", "", cd_so_bo, flags=re.IGNORECASE)
        cd_sau_mo = cd_clean.split("-")[0].strip() or cd_so_bo

    # 4. Trích xuất Ngày hậu phẫu
    ngay_hp_val = str(st.session_state.get("ngay_hau_phau", "")).strip()
    m_hp = re.search(r"(?:ngày\s*(?:thứ)?\s*)(\d+)", ngay_hp_val, re.IGNORECASE)
    if m_hp:
        ngay_hp_str = f"ngày thứ {m_hp.group(1)}"
    elif ngay_hp_val:
        ngay_hp_str = ngay_hp_val
    else:
        ngay_hp_str = "ngày thứ ..."

    return (
        f"Bệnh nhân {gioi_tinh.lower()} {tuoi_str} vào viện vì {ly_do}, "
        f"chẩn đoán trước mổ là {cd_truoc_mo}, được mổ bằng phương pháp {pp_mo}, "
        f"mổ {loai_mo}, chẩn đoán sau mổ là {cd_sau_mo}. "
        f"Quá trình mổ không có biến chứng. Hiện tại hậu phẫu {ngay_hp_str}. "
        f"Qua thăm khám và hỏi bệnh phát hiện các hội chứng và triệu chứng sau:"
    )

def ui_tom_tat(num):
    col_tt_title, col_tt_btn = st.columns([1.5, 0.5])
    with col_tt_title:
        st.markdown(f"**{num}. Tóm tắt bệnh án:**")
    with col_tt_btn:
        # Tự động chọn câu dẫn phù hợp với loại bệnh án
        help_text = "Tự động trích xuất thông tin phẫu thuật và điền câu dẫn" if loai_benh_an == "Hậu phẫu" else "Tự động tính ngày và điền câu dẫn mở đầu"
        if st.button("⚡ Tạo câu dẫn", key="btn_auto_cau_dan_tt", help=help_text, use_container_width=True):
            if loai_benh_an == "Hậu phẫu":
                cau_dan_moi = generate_intro_tom_tat_hau_phau()
            else:
                cau_dan_moi = generate_intro_tom_tat_noi_khoa()

            current_tt = str(st.session_state.get("tom_tat", "")).strip()
            if current_tt:
                lines = current_tt.split("\n")
                if len(lines) > 1 and (lines[1].strip().startswith("-") or lines[1].strip().startswith("*")):
                    st.session_state["tom_tat"] = cau_dan_moi + "\n" + "\n".join(lines[1:])
                else:
                    st.session_state["tom_tat"] = cau_dan_moi + "\n" + current_tt
            else:
                st.session_state["tom_tat"] = cau_dan_moi
            st.rerun()

    st.text_area(
        f"{num}. Tóm tắt bệnh án:", 
        key="tom_tat", 
        height=130, 
        label_visibility="collapsed",
        placeholder="Dòng đầu tiên là câu dẫn tóm tắt. Các dòng tiếp theo ghi các hội chứng và triệu chứng có giá trị..."
    )

def ui_cdsb(num_sb, num_pb, num_bl):
    # Hàng 1: Tiêu đề + Nút bấm căn ngang hàng
    col_h_left, col_h_right_title, col_h_right_btn = st.columns([1, 0.65, 0.35])
    with col_h_left:
        st.markdown(f"**{num_sb}. Chẩn đoán sơ bộ:**")
    with col_h_right_title:
        st.markdown(f"**{num_pb}. Chẩn đoán phân biệt:**")
    with col_h_right_btn:
        btn_ai_cdpb = st.button("🪄 Làm phép", key="btn_ai_cdpb", type="primary", use_container_width=True)

    # Xử lý logic AI khi bấm nút
    if btn_ai_cdpb:
        if "GEMINI_API_KEY" not in st.secrets:
            st.error("⚠️ Chưa cài đặt API Key bí mật!")
        elif not str(st.session_state.get("chan_doan_so_bo", "")).strip():
            st.warning("⚠️ Vui lòng nhập Chẩn đoán sơ bộ!")
        else:
            with st.spinner("AI đang phân tích lập luận lâm sàng để tạo Chẩn đoán phân biệt & Biện luận..."):
                try:
                    benh_su_str = get_benh_su_text_for_ai()
                    context_cdpb = (
                        f"Loại bệnh án: {loai_benh_an}\n"
                        f"Bệnh nhân: {st.session_state.get('tuoi')} tuổi, Giới tính: {st.session_state.get('gioi_tinh')}\n"
                        f"Lý do vào viện: {st.session_state.get('ly_do_vao_vien')}\n"
                        f"Bệnh sử: {benh_su_str}\n"
                        f"Tiền sử: {st.session_state.get('ts_noi_khoa')} | Ngoại khoa: {st.session_state.get('ts_ngoai_khoa')}\n"
                        f"Sinh hiệu: Mạch {st.session_state.get('sh_mach')}, HA {st.session_state.get('sh_ha')}, Nhiệt độ {st.session_state.get('sh_nhiet_do')}\n"
                        f"Khám toàn thân: {st.session_state.get('kham_toan_than')}\n"
                    )
                    if loai_benh_an == "Hậu phẫu":
                        context_cdpb += (
                            f"Ngày hậu phẫu: {st.session_state.get('ngay_hau_phau')}\n"
                            f"Khám vết mổ: {st.session_state.get('kham_vet_mo')}\n"
                            f"Khám dẫn lưu: {st.session_state.get('kham_dan_luu')}\n"
                        )
                    context_cdpb += f"CHẨN ĐOÁN SƠ BỘ: {st.session_state.get('chan_doan_so_bo')}"

                    model = get_feature_model("KEY_AI", "gemini-3.1-flash-lite")
                    prompt_cdpb = f"""
                    Bạn là một bác sĩ lâm sàng thực thụ và giàu kinh nghiệm. Hãy nhìn vào toàn thể ca bệnh dưới đây, phân tích logic giữa bệnh cảnh, triệu chứng cơ năng, thực thể và chẩn đoán sơ bộ để đưa ra:
                    1. Danh sách CHẨN ĐOÁN PHÂN BIỆT (Differential Diagnosis): sắp xếp thứ tự từ khả năng cao nhất đến thấp hơn, từ bệnh lý cấp cứu nguy hiểm đến ít cấp cứu hơn.
                    2. BIỆN LUẬN CHẨN ĐOÁN SƠ BỘ: Lập luận chặt chẽ vì sao nghĩ đến chẩn đoán sơ bộ và vì sao cần phân biệt với các bệnh lý nêu trên.

                    Dữ kiện ca bệnh:
                    {context_cdpb}

                    YÊU CẦU ĐẦU RA (Xuất ra đúng 2 khối nhãn sau, không viết thêm lời dẫn chào hỏi):
                    [CHAN_DOAN_PHAN_BIET]
                    1. Tên bệnh A
                    2. Tên bệnh B
                    3. Tên bệnh C

                    [BIEN_LUAN_SO_BO]
                    (Nội dung đoạn văn biện luận logic, súc tích).
                    """
                    res_text = model.generate_content(prompt_cdpb).text

                    if "[CHAN_DOAN_PHAN_BIET]" in res_text and "[BIEN_LUAN_SO_BO]" in res_text:
                        parts = res_text.split("[BIEN_LUAN_SO_BO]")
                        st.session_state["chan_doan_phan_biet"] = parts[0].replace("[CHAN_DOAN_PHAN_BIET]", "").strip()
                        st.session_state["bien_luan"] = parts[1].strip()
                        st.toast("✨ Đã tạo gợi ý Chẩn đoán phân biệt & Biện luận!", icon="🪄")
                        st.rerun()
                    else:
                        st.session_state["chan_doan_phan_biet"] = res_text.strip()
                        st.rerun()
                except Exception as e:
                    st.error(f"Lỗi AI: {e}")

    # Hàng 2: Hai ô nhập liệu ngang hàng nhau, cùng chiều cao
    c_cd1, c_cd2 = st.columns(2)
    with c_cd1:
        placeholder_cd = "Hậu phẫu ngày thứ [X]... mổ phiên/cấp cứu do [Bệnh lý]..." if loai_benh_an == "Hậu phẫu" else "Chẩn đoán sơ bộ..."
        st.text_area(f"{num_sb}. Chẩn đoán sơ bộ:", key="chan_doan_so_bo", height=100, placeholder=placeholder_cd, label_visibility="collapsed")
    with c_cd2:
        st.text_area(f"{num_pb}. Chẩn đoán phân biệt:", key="chan_doan_phan_biet", height=100, label_visibility="collapsed")

    # Hàng 3: Biện luận chẩn đoán sơ bộ
    st.markdown(f"**{num_bl}. Biện luận chẩn đoán sơ bộ:**")
    st.text_area(f"{num_bl}. Biện luận chẩn đoán sơ bộ:", key="bien_luan", height=130, label_visibility="collapsed")

def ui_cls(num_dx, num_kq):
    st.markdown(f"<div class='sub-section-header'>{num_dx}. Đề xuất cận lâm sàng</div>", unsafe_allow_html=True)
    if st.button("🪄 Làm phép", type="primary", key="btn_ai_cls"):
        if "GEMINI_API_KEY" not in st.secrets: st.error("⚠️ Chưa cài đặt API Key!")
        else:
            with st.spinner("AI đang phân tích chỉ định cận lâm sàng tối ưu..."):
                try:
                    benh_su_str = get_benh_su_text_for_ai()
                    context_cls = f"Loại bệnh án: {loai_benh_an}\nTuổi: {st.session_state.get('tuoi')}, Giới tính: {st.session_state.get('gioi_tinh')}\nBệnh sử: {benh_su_str}\nChẩn đoán sơ bộ: {st.session_state.get('chan_doan_so_bo')}"
                    model = get_feature_model("KEY_AI", "gemini-3.1-flash-lite")
                    prompt_cls = f"Bạn là bác sĩ lâm sàng. Phân tích ca bệnh ({context_cls}) để chỉ định CẬN LÂM SÀNG. Nếu là Hậu phẫu, ưu tiên các xét nghiệm theo dõi biến chứng mổ. Trả về đúng 3 nhãn: [CLS_XAC_DINH], [CLS_DIEU_TRI], [CLS_KHAC] dưới dạng xuống dòng, không dùng gạch đầu dòng."
                    res_cls_text = model.generate_content(prompt_cls).text
                    if "[CLS_XAC_DINH]" in res_cls_text and "[CLS_DIEU_TRI]" in res_cls_text:
                        p1 = res_cls_text.split("[CLS_DIEU_TRI]")
                        part_xd = p1[0].replace("[CLS_XAC_DINH]", "").strip()
                        if "[CLS_KHAC]" in p1[1]:
                            p2 = p1[1].split("[CLS_KHAC]")
                            part_dt, part_khac = p2[0].strip(), p2[1].strip()
                        else:
                            part_dt, part_khac = p1[1].strip(), ""
                        st.session_state["cls_dx_xac_dinh"] = part_xd
                        st.session_state["cls_dx_dieu_tri"] = part_dt
                        st.session_state["cls_dx_khac"] = part_khac
                        st.success("✨ Đã gợi ý danh mục CLS thành công!")
                        st.rerun()
                    else: st.error("AI trả về sai định dạng.")
                except Exception as e: st.error(f"Lỗi AI: {e}")

    c_cls1, c_cls2, c_cls3 = st.columns(3)
    with c_cls1: st.text_area("1. Phục vụ chẩn đoán xác định:", key="cls_dx_xac_dinh", height=130)
    with c_cls2: st.text_area("2. Phục vụ điều trị:", key="cls_dx_dieu_tri", height=130)
    with c_cls3: st.text_area("3. Cận lâm sàng khác:", key="cls_dx_khac", height=130)
        
    st.markdown(f"<div class='sub-section-header'>{num_kq}. Cận lâm sàng đã có (Hiện có {st.session_state['so_hang_cls']} hàng)</div>", unsafe_allow_html=True)
    with st.container():
        col_ocr_file, col_ocr_act = st.columns([2.5, 1])
        with col_ocr_file: lab_photos = st.file_uploader("📷 Tải lên ảnh phiếu xét nghiệm (cho phép nhiều ảnh):", type=["png", "jpg", "jpeg"], accept_multiple_files=True, key="uploader_ocr_lab_multi")
        with col_ocr_act: 
            st.write(""); st.write("")
            btn_ocr = st.button("⚡ Phân tích tất cả ảnh", type="primary", use_container_width=True, key="btn_ocr_lab_batch")

        if btn_ocr and lab_photos:
            vision_model = get_feature_model("KEY_OCR", "gemini-3.1-flash-lite")
            if not vision_model: st.error("⚠️ Hệ thống chưa được cấu hình API Key!")
            else:
                progress_bar = st.progress(0, text="Bắt đầu phân tích...")
                ocr_prompt = "Bạn là bác sĩ xét nghiệm. Đọc phiếu này và trả về JSON có 2 khóa: 'ket_qua' (liệt kê chỉ số dạng \n- ) và 'phien_giai' (biện luận chỉ số bất thường)."
                so_hang_cls = int(st.session_state.get("so_hang_cls", 3))
                last_used_idx = -1
                for r in range(so_hang_cls):
                    val_k, val_p = str(st.session_state.get(f"cls_kq_{r}", "")).strip(), str(st.session_state.get(f"cls_pg_{r}", "")).strip()
                    if (val_k and val_k not in ["None", "-"]) or (val_p and val_p not in ["None", "-"]): last_used_idx = r

                start_row = last_used_idx + 1
                empty_rows = [start_row + i for i in range(len(lab_photos))]
                if start_row + len(lab_photos) > so_hang_cls: st.session_state["so_hang_cls"] = start_row + len(lab_photos)

                thanh_cong = 0
                for idx, photo in enumerate(lab_photos):
                    target_row = empty_rows[idx]
                    progress_bar.progress(int((idx + 1) / len(lab_photos) * 100), text=f"Xử lý ảnh {idx + 1}/{len(lab_photos)}...")
                    try:
                        img_input = optimize_lab_image(photo)
                        resp = vision_model.generate_content([ocr_prompt, img_input])
                        raw_text = resp.text.strip()
                        if raw_text.startswith("```json"): raw_text = raw_text[7:]
                        elif raw_text.startswith("```"): raw_text = raw_text[3:]
                        if raw_text.endswith("```"): raw_text = raw_text[:-3]

                        lab_data = json.loads(raw_text.strip())
                        if isinstance(lab_data, dict):
                            st.session_state[f"cls_kq_{target_row}"] = lab_data.get("ket_qua", "")
                            st.session_state[f"cls_pg_{target_row}"] = lab_data.get("phien_giai", "")
                            thanh_cong += 1
                    except Exception as e: st.warning(f"Không thể phân tích ảnh {photo.name}: {e}")

                progress_bar.empty()
                if thanh_cong > 0:
                    st.toast(f"✅ Đã phân tích xong {thanh_cong} ảnh!", icon="🧪")
                    st.rerun()
        st.divider()

    for i in range(st.session_state["so_hang_cls"]):
        st.markdown(f"**Hàng {i + 1}:**")
        col_left, col_right = st.columns([1, 1])
        with col_left:
            st.text_area(f"Kết quả cận lâm sàng {i + 1}:", key=f"cls_kq_{i}", height=75)
            img = st.file_uploader(f"Đính kèm ảnh cho hàng {i + 1}:", type=["png", "jpg", "jpeg"], key=f"uploader_cls_img_{i}")
            if img:
                uploaded_imgs[f"cls_img_{i}"] = img
                st.image(img, width=180, caption=f"Ảnh hàng {i + 1}")
        with col_right:
            st.text_area(f"Biện giải cận lâm sàng {i + 1}:", key=f"cls_pg_{i}", height=130)
        st.divider()

    col_btn_them, col_btn_bot, _ = st.columns([2, 2, 6])
    with col_btn_them:
        if st.button("➕ Thêm kết quả cận lâm sàng"): st.session_state["so_hang_cls"] += 1; st.rerun()
    with col_btn_bot:
        if st.session_state["so_hang_cls"] > 1:
            if st.button("➖ Bớt hàng cuối"): st.session_state["so_hang_cls"] -= 1; st.rerun()

def ui_cdxd(num_xd, num_blxd):
    placeholder_xd = "Phẫu thuật [Tên PT] mổ [phiên/cấp cứu] ngày thứ [X] do [Bệnh lý] hiện tại [ổn định/biến chứng...]" if loai_benh_an == "Hậu phẫu" else "Chẩn đoán xác định..."
    st.text_area(f"{num_xd}. Chẩn đoán xác định:", key="chan_doan_xac_dinh", height=90, placeholder=placeholder_xd)
    st.text_area(f"{num_blxd}. Biện luận chẩn đoán xác định:", key="bien_luan_xac_dinh", height=110)

tab1, tab2, tab3 = st.tabs(["Nhập liệu hồ sơ", "Xuất tập tin", "Phản biện lâm sàng"])

with tab1:
    with st.expander("I. PHẦN HÀNH CHÍNH", expanded=True):
        c_hc1, c_hc2, c_hc3 = st.columns(3)
        with c_hc1:
            st.text_input("Họ và tên người bệnh", key="ho_ten", placeholder="Nguyễn Văn A")
            st.text_input("Dân tộc", key="dan_tok", placeholder="Kinh, Tày, Nùng...")
        with c_hc2:
            st.number_input("Tuổi", min_value=0, max_value=120, key="tuoi")
            st.text_input("Nghề nghiệp", key="nghe_nghiep", placeholder="Kỹ sư, Hưu trí, Nông dân...")
        with c_hc3:
            st.selectbox("Giới tính", ["Nam", "Nữ", "Khác"], key="gioi_tinh")
            st.text_input("Khoa / Phòng điều trị", key="khoa_phong", placeholder="Khoa Ngoại Tiêu Hóa, Nội Tim mạch...")
        
        c_hc4, c_hc5, c_hc6 = st.columns(3)
        with c_hc4: st.text_input("Địa chỉ", key="dia_chi", placeholder="Quận Đống Đa, TP. Hà Nội")
        with c_hc5: st.text_input("Bác sĩ hoặc Sinh viên phụ trách", key="sinh_vien", placeholder="Bác sĩ nội trú, Sinh viên Y...")
        with c_hc6: st.text_input("Ngày giờ vào viện", key="ngay_vao_vien")

    with st.expander("II VÀ III. LÝ DO VÀO VIỆN VÀ BỆNH SỬ", expanded=True):
        st.text_area("Lý do vào viện:", key="ly_do_vao_vien", placeholder="Ví dụ: Giống bệnh án tiền phẫu", height=65)
        
        if loai_benh_an == "Hậu phẫu":
            st.markdown("**BỆNH SỬ HẬU PHẪU:**")
            st.text_area("1. Tình trạng trước mổ:", key="bs_truoc_mo", height=90, placeholder="Chỉ nêu các triệu chứng chính và Chẩn đoán trước mổ...")
            st.text_area("2. Tình trạng trong mổ:", key="bs_trong_mo", height=90, placeholder="Mổ phiên hay cấp cứu, ngày giờ mổ, vô cảm, phẫu thuật, tổn thương, tai biến trong mổ (nếu có)...")
            st.text_area("3. Quá trình sau mổ:", key="bs_sau_mo", height=90, placeholder="Từ lúc rời phòng hồi tỉnh đến nay: Tri giác, đau, trung tiện, tiểu tiện, tình trạng dẫn lưu, ăn uống...")
        else:
            st.text_area("Bệnh sử:", key="benh_su", placeholder="Mô tả hoàn cảnh khởi phát, triệu chứng cơ năng điển hình...", height=130)

    with st.expander("IV. TIỀN SỬ", expanded=True):
        c_ts1, c_ts2 = st.columns(2)
        with c_ts1:
            st.markdown("<div class='sub-section-header'>1. Tiền sử nội khoa</div>", unsafe_allow_html=True)
            st.text_area("Nội dung tiền sử nội khoa:", key="ts_noi_khoa", height=90, label_visibility="collapsed")
            st.markdown("<div class='sub-section-header'>2. Tiền sử ngoại khoa và dị ứng</div>", unsafe_allow_html=True)
            st.text_area("Nội dung tiền sử ngoại khoa và dị ứng:", key="ts_ngoai_khoa", height=90, label_visibility="collapsed")
        with c_ts2:
            st.markdown("<div class='sub-section-header'>3. Lối sống và thói quen</div>", unsafe_allow_html=True)
            st.text_area("Nội dung lối sống và thói quen:", key="ts_loi_song", height=90, label_visibility="collapsed")
            st.markdown("<div class='sub-section-header'>4. Tiền sử gia đình</div>", unsafe_allow_html=True)
            st.text_area("Nội dung tiền sử gia đình:", key="ts_gia_dinh", height=90, label_visibility="collapsed")

    with st.expander("V. THĂM KHÁM LÂM SÀNG", expanded=True):
        # Không hiển thị mục "Khám vào viện" nếu là Hậu phẫu
        if loai_benh_an == "Hậu phẫu":
            st.markdown("<div class='sub-section-header'>1. Thăm khám hiện tại - Toàn thân & Sinh hiệu</div>", unsafe_allow_html=True)
            st.text_input("Khám hậu phẫu ngày thứ mấy? Giờ thứ mấy?", key="ngay_hau_phau", placeholder="VD: Ngày thứ 3 sau mổ (Giờ thứ 72)...")
        else:
            st.markdown("<div class='sub-section-header'>1. Thăm khám lúc vào viện</div>", unsafe_allow_html=True)
            st.text_area("Nội dung khám lúc vào viện:", key="kham_vao_vien", height=80, label_visibility="collapsed")
            st.markdown("<div class='sub-section-header'>2. Thăm khám hiện tại - Toàn thân & Sinh hiệu</div>", unsafe_allow_html=True)

        col_tt_mo_ta, col_tt_sh = st.columns([1.2, 1])
        with col_tt_mo_ta:
            st.markdown("**Mô tả khám toàn thân:**")
            st.text_area("Nội dung khám toàn thân:", key="kham_toan_than", height=175, label_visibility="collapsed", placeholder="- Tri giác, tiếp xúc (tỉnh/mê, GCS...)\n- Da niêm mạc (hồng, nhợt, vàng da, xuất huyết dưới da...)\n- Lông tóc móng, tuyến giáp, hạch ngoại vi, phù...")
        
        with col_tt_sh:
            st.markdown("**Dấu hiệu sinh tồn (Vital Signs):**")
            c_sh1, c_sh2 = st.columns(2)
            with c_sh1:
                st.text_input("Mạch (lần/phút):", key="sh_mach", placeholder="VD: 80")
                st.text_input("Huyết áp (mmHg):", key="sh_ha", placeholder="VD: 120/80")
                st.number_input("Cân nặng (kg):", key="sh_can_nang", min_value=0.0, max_value=250.0, value=float(st.session_state.get("sh_can_nang", 0.0)), step=0.5)
            with c_sh2:
                st.text_input("Nhiệt độ (°C):", key="sh_nhiet_do", placeholder="VD: 37.0")
                st.text_input("Nhịp thở (lần/phút):", key="sh_nhip_tho", placeholder="VD: 18")
                st.number_input("Chiều cao (cm):", key="sh_chieu_cao", min_value=0.0, max_value=230.0, value=float(st.session_state.get("sh_chieu_cao", 0.0)), step=1.0)
            
            w, h = st.session_state.get("sh_can_nang", 0.0), st.session_state.get("sh_chieu_cao", 0.0)
            if w > 0 and h > 0:
                bmi_val = round(w / ((h / 100.0) ** 2), 1)
                bmi_eval = "Gầy / Thiếu cân" if bmi_val < 18.5 else "Bình thường" if bmi_val <= 22.9 else "Thừa cân" if bmi_val <= 24.9 else "Béo phì"
                st.session_state["sh_bmi"], st.session_state["sh_bmi_eval"] = str(bmi_val), bmi_eval
                st.caption(f"📊 **BMI:** `{bmi_val} kg/m²` — **Đánh giá:** *{bmi_eval}*")
            else:
                st.session_state["sh_bmi"], st.session_state["sh_bmi_eval"] = "", ""
        
        if loai_benh_an == "Hậu phẫu":
            st.markdown("<div class='sub-section-header'>2. Thăm khám Vết mổ & Dẫn lưu</div>", unsafe_allow_html=True)
            c_vm, c_dl = st.columns(2)
            with c_vm: st.text_area("Tình trạng vết mổ:", key="kham_vet_mo", height=85, placeholder="Ví dụ: Vết mổ khô, không sưng đỏ, chân chỉ không nề...")
            with c_dl: st.text_area("Tình trạng ống dẫn lưu:", key="kham_dan_luu", height=85, placeholder="Ví dụ: Dẫn lưu ổ bụng ra 20ml dịch hồng nhạt...")
            st.markdown("<div class='sub-section-header'>3. Thăm khám hiện tại - Các cơ quan</div>", unsafe_allow_html=True)
        else:
            st.markdown("<div class='sub-section-header'>3. Thăm khám hiện tại - Các cơ quan</div>", unsafe_allow_html=True)

        def xu_ly_dien_kham_binh_thuong():
            dem = 0
            for k_cq, norm_val in NORMAL_ORGAN_FINDINGS.items():
                noi_dung = str(st.session_state.get(k_cq, "") or "").strip()
                if not noi_dung:
                    st.session_state[k_cq] = norm_val
                    dem += 1
            st.session_state["_msg_dien_cq"] = dem

        def xu_ly_xoa_cac_co_quan():
            for k_cq in NORMAL_ORGAN_FINDINGS.keys():
                st.session_state[k_cq] = ""
            st.session_state["_msg_xoa_cq"] = True

        col_btn_fill, col_clear_cq = st.columns([2, 1])
        with col_btn_fill:
            st.button("⚡ Điền khám bình thường cho các cơ quan để trống", on_click=xu_ly_dien_kham_binh_thuong, use_container_width=True)
        with col_clear_cq:
            st.button("🔄 Đặt lại các cơ quan", on_click=xu_ly_xoa_cac_co_quan, use_container_width=True)

        if "_msg_dien_cq" in st.session_state:
            d = st.session_state.pop("_msg_dien_cq")
            if d > 0:
                st.toast(f"Đã điền mẫu cho {d} cơ quan còn lại!", icon="✨")
            else:
                st.info("Tất cả các cơ quan đều đã có dữ liệu.")

        if st.session_state.pop("_msg_xoa_cq", False):
            st.toast("Đã làm trống các ô khám cơ quan!", icon="🧹")
            
        st.markdown("---")

        ORGAN_DEF = [{"key": "kham_tuan_hoan", "name": "Tuần hoàn"}, {"key": "kham_ho_hap", "name": "Hô hấp"}, {"key": "kham_tieu_hoa", "name": "Tiêu hóa"}, {"key": "kham_than_kinh", "name": "Thần kinh"}, {"key": "kham_tiet_nieu", "name": "Thận - Tiết niệu"}, {"key": "kham_co_xuong_khop", "name": "Cơ xương khớp"}, {"key": "kham_co_quan_khac", "name": "Các cơ quan khác"}]
        selected_organ_name = st.selectbox("Chọn cơ quan chuyên khoa ưu tiên:", ["Không ưu tiên (Thứ tự mặc định)"] + [item["name"] for item in ORGAN_DEF], index=0, key="uu_tien_co_quan")

        if selected_organ_name != "Không ưu tiên (Thứ tự mặc định)":
            fav = next(item for item in ORGAN_DEF if item["name"] == selected_organ_name)
            others = [item for item in ORGAN_DEF if item["name"] != selected_organ_name]
            st.markdown(f"**{fav['name'].upper()} (CƠ QUAN CHUYÊN KHOA TRỌNG ĐIỂM):**")
            st.text_area(f"Khám {fav['name']}:", key=fav["key"], height=130)
            st.markdown("---")
            st.markdown("**Các cơ quan khác:**")
            c_cq1, c_cq2 = st.columns(2)
            half = len(others) // 2 + len(others) % 2
            with c_cq1:
                for idx, org in enumerate(others[:half]): st.text_area(f"{org['name']}:", key=org["key"], height=85)
            with c_cq2:
                for idx, org in enumerate(others[half:]): st.text_area(f"{org['name']}:", key=org["key"], height=85)
        else:
            c_cq1, c_cq2 = st.columns(2)
            with c_cq1:
                st.text_area("Tuần hoàn:", key="kham_tuan_hoan", height=85)
                st.text_area("Hô hấp:", key="kham_ho_hap", height=85)
                st.text_area("Tiêu hóa:", key="kham_tieu_hoa", height=85)
                st.text_area("Thần kinh:", key="kham_than_kinh", height=85)
            with c_cq2:
                st.text_area("Thận - Tiết niệu:", key="kham_tiet_nieu", height=85)
                st.text_area("Cơ xương khớp:", key="kham_co_xuong_khop", height=85)
                st.text_area("Các cơ quan khác:", key="kham_co_quan_khac", height=85)

    # --- KHỐI ĐỘNG CHUYỂN MẠCH VỊ TRÍ THEO LOẠI BỆNH ÁN ---
    if loai_benh_an == "Hậu phẫu":
        with st.expander("VI VÀ VII. CHẨN ĐOÁN SƠ BỘ VÀ PHÂN BIỆT", expanded=True):
            ui_cdsb("VI", "VII", "VIII")
        with st.expander("IX VÀ X. CẬN LÂM SÀNG", expanded=True):
            ui_cls("IX", "X")
        with st.expander("XI VÀ XII. TÓM TẮT BỆNH ÁN VÀ CHẨN ĐOÁN XÁC ĐỊNH", expanded=True):
            ui_tom_tat("XI")
            ui_cdxd("XII", "XIII")
    else:
        with st.expander("VI ĐẾN IX. TÓM TẮT VÀ BIỆN LUẬN CHẨN ĐOÁN SƠ BỘ", expanded=True):
            ui_tom_tat("VI")
            ui_cdsb("VII", "VIII", "IX")
        with st.expander("X VÀ XI. CẬN LÂM SÀNG", expanded=True):
            ui_cls("X", "XI")
        with st.expander("XII VÀ XIII. CHẨN ĐOÁN XÁC ĐỊNH VÀ BIỆN LUẬN", expanded=True):
            ui_cdxd("XII", "XIII")

    with st.expander("XIV. HƯỚNG DẪN VÀ KẾ HOẠCH ĐIỀU TRỊ", expanded=True):
        if st.button("🪄 Làm phép", key="btn_ai_dt", type="primary"):
            if "GEMINI_API_KEY" not in st.secrets: st.error("⚠️ Chưa cài đặt API Key!")
            else:
                with st.spinner("AI đang phân tích phác đồ điều trị..."):
                    try:
                        cls_da_co_str = "".join([f"+ {st.session_state.get(f'cls_kq_{i}', '')} -> {st.session_state.get(f'cls_pg_{i}', '')}\n" for i in range(st.session_state.get("so_hang_cls", 3)) if st.session_state.get(f'cls_kq_{i}', '').strip()])
                        context_dt = f"Loại: {loai_benh_an}\nBệnh nhân: {st.session_state.get('tuoi')} tuổi, {st.session_state.get('gioi_tinh')}\nTiền sử: {st.session_state.get('ts_noi_khoa')}\nChẩn đoán: {st.session_state.get('chan_doan_xac_dinh')}\nCLS quan trọng:\n{cls_da_co_str}"
                        model = get_feature_model("KEY_AI", "gemini-3.1-flash-lite")
                        prompt_dt = f"Bạn là bác sĩ điều trị. Xây dựng phác đồ cho ca bệnh ({context_dt}). Yêu cầu trả về đúng 3 tag: [MUC_TIEU], [DIEU_TRI_CU_THE] (ghi rõ thuốc/chăm sóc vết mổ nếu hậu phẫu), [THEO_DOI] theo định dạng text trơn không gạch đầu dòng."
                        txt = model.generate_content(prompt_dt).text

                        if "[MUC_TIEU]" in txt and "[DIEU_TRI_CU_THE]" in txt and "[THEO_DOI]" in txt:
                            p1 = txt.split("[DIEU_TRI_CU_THE]")
                            st.session_state["dt_muc_tieu"] = p1[0].replace("[MUC_TIEU]", "").strip()
                            p2 = p1[1].split("[THEO_DOI]")
                            st.session_state["dt_cu_the"] = p2[0].strip()
                            st.session_state["dt_theo_doi"] = p2[1].strip()
                            st.success("✨ Đã lên phác đồ điều trị thành công!")
                            st.rerun()
                        else: st.error("AI phản hồi sai cấu trúc.")
                    except Exception as e: st.error(f"Lỗi AI: {e}")

        c_mt, c_ct, c_td = st.columns(3)
        with c_mt: st.text_area("1. Mục tiêu điều trị:", key="dt_muc_tieu", height=220)
        with c_ct: st.text_area("2. Điều trị cụ thể:", key="dt_cu_the", height=220)
        with c_td: st.text_area("3. Theo dõi:", key="dt_theo_doi", height=220)

    with st.expander("XV VÀ XVI. TIÊN LƯỢNG VÀ TƯ VẤN", expanded=True):
        if st.button("🪄 Làm phép", type="primary", key="btn_ai_tienluong"):
            if "GEMINI_API_KEY" not in st.secrets: st.error("⚠️ Chưa cài đặt API Key!")
            else:
                with st.spinner("AI đang phân tích logic lâm sàng..."):
                    try:
                        context = f"Loại: {loai_benh_an}\nTuổi: {st.session_state.get('tuoi')}, Giới tính: {st.session_state.get('gioi_tinh')}\nChẩn đoán: {st.session_state.get('chan_doan_xac_dinh')}\nĐiều trị: {st.session_state.get('dt_cu_the')}"
                        model = get_feature_model("KEY_AI", "gemini-3.1-flash-lite")
                        prompt = f"Bạn là bác sĩ lâm sàng. Đưa ra TIÊN LƯỢNG và TƯ VẤN cho ca bệnh ({context}). Yêu cầu trả về đúng 2 tag: [TIEN_LUONG] và [TU_VAN]."
                        res_text = model.generate_content(prompt).text

                        if "[TIEN_LUONG]" in res_text and "[TU_VAN]" in res_text:
                            parts = res_text.split("[TU_VAN]")
                            st.session_state["tien_luong"] = parts[0].replace("[TIEN_LUONG]", "").strip()
                            st.session_state["tu_van"] = parts[1].strip()
                            st.success("✨ Đã tạo gợi ý thành công!")
                            st.rerun() 
                        else: st.error("AI trả về sai định dạng.")
                    except Exception as e: st.error(f"Lỗi AI: {e}")

        c_pl, c_tv = st.columns(2)
        with c_pl: st.text_area("XV. Tiên lượng:", key="tien_luong", height=250)
        with c_tv: st.text_area("XVI. Tư vấn:", key="tu_van", height=250)

# Gom dữ liệu để xuất file
data_benh_an = {k: st.session_state.get(k, "") for k in FIELDS_TO_SAVE}
data_benh_an["loai_benh_an"] = loai_benh_an
data_benh_an["sh_mach"] = str(st.session_state.get("sh_mach", "")).strip()
data_benh_an["sh_nhiet_do"] = str(st.session_state.get("sh_nhiet_do", "")).strip()
data_benh_an["sh_ha"] = str(st.session_state.get("sh_ha", "")).strip()
data_benh_an["sh_nhip_tho"] = str(st.session_state.get("sh_nhip_tho", "")).strip()
data_benh_an["sh_can_nang"] = str(st.session_state.get("sh_can_nang", 0.0))
data_benh_an["sh_chieu_cao"] = str(st.session_state.get("sh_chieu_cao", 0.0))
data_benh_an["sh_bmi"] = str(st.session_state.get("sh_bmi", ""))
data_benh_an["sh_bmi_eval"] = str(st.session_state.get("sh_bmi_eval", ""))
data_benh_an["so_hang_cls"] = st.session_state.get("so_hang_cls", 3)
for i in range(data_benh_an["so_hang_cls"]):
    data_benh_an[f"cls_kq_{i}"] = st.session_state.get(f"cls_kq_{i}", "")
    data_benh_an[f"cls_pg_{i}"] = st.session_state.get(f"cls_pg_{i}", "")
if 'uploaded_imgs' in locals(): data_benh_an.update(uploaded_imgs)

# --- TAB 2: XEM TRƯỚC VÀ XUẤT TẬP TIN ---
with tab2:
    st.markdown("<div class='sidebar-header-amboss'>XEM TRƯỚC THÔNG TIN TỔNG QUAN</div>", unsafe_allow_html=True)
    ho_ten_val = str(st.session_state.get("ho_ten", "")).strip()
    if ho_ten_val:
        st.info(f"Bệnh nhân: {ho_ten_val.upper()} | {st.session_state.get('tuoi')} tuổi | Giới tính: {st.session_state.get('gioi_tinh')} | Loại Bệnh Án: {loai_benh_an}")
        st.write(f"Khoa phòng: {st.session_state.get('khoa_phong', 'Chưa điền')} | Lý do vào viện: {st.session_state.get('ly_do_vao_vien')}")
        if st.session_state.get("chan_doan_so_bo"): st.write(f"Chẩn đoán sơ bộ: {st.session_state.get('chan_doan_so_bo')}")
        if st.session_state.get("chan_doan_xac_dinh"): st.markdown(f"<div class='highlight-dx'>Chẩn đoán xác định: {st.session_state.get('chan_doan_xac_dinh')}</div>", unsafe_allow_html=True)
        so_hang = st.session_state.get("so_hang_cls", 3)
        dem_cls = sum(1 for i in range(so_hang) if str(st.session_state.get(f"cls_kq_{i}", "")).strip() or (locals().get('uploaded_imgs') and uploaded_imgs.get(f"cls_img_{i}")))
        st.write(f"Số lượng cận lâm sàng đã nhập vào bảng: {dem_cls}/{so_hang} hàng")
    else:
        st.warning("Vui lòng điền thông tin bên tab Nhập liệu hồ sơ.")

    st.markdown("---")
    col_dl_pdf, col_dl_pptx = st.columns(2)
    with col_dl_pdf:
        if st.button("📄 Tạo & Xem trước tập tin PDF", type="primary", use_container_width=True):
            if not ho_ten_val: st.error("Vui lòng điền tối thiểu Họ và tên người bệnh trước khi xuất tập tin!")
            elif not os.path.exists("Roboto-Regular.ttf") or not os.path.exists("Roboto-Bold.ttf"): st.error("Chưa tìm thấy tập tin font 'Roboto-Regular.ttf' và 'Roboto-Bold.ttf' trong cùng thư mục với app.py!")
            else:
                with st.spinner("Đang kết xuất văn bản PDF..."):
                    pdf_bytes = export_pdf(data_benh_an)
                    st.session_state["pdf_bytes_preview"] = pdf_bytes
                    st.session_state["ten_file_pdf"] = f"Benh_an_{'Hau_phau_' if loai_benh_an == 'Hậu phẫu' else ''}{ho_ten_val.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pdf"

        if st.session_state.get("pdf_bytes_preview"):
            st.download_button("📥 Tải PDF về máy", data=st.session_state["pdf_bytes_preview"], file_name=st.session_state.get("ten_file_pdf", "benh_an.pdf"), mime="application/pdf", use_container_width=True)

    with col_dl_pptx:
        if st.button("Tạo tập tin PowerPoint (PPTX)", type="secondary", use_container_width=True):
            if not ho_ten_val: st.error("Vui lòng điền tối thiểu Họ và tên người bệnh!")
            else:
                with st.spinner("Đang kết xuất bản trình chiếu PowerPoint..."):
                    pptx_bytes = export_pptx(data_benh_an)
                    ten_file_pptx = f"Trinh_chieu_Benh_an_{'Hau_phau_' if loai_benh_an == 'Hậu phẫu' else ''}{ho_ten_val.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx"
                    st.success("Tạo PowerPoint thành công!")
                    st.download_button("Nhấn vào đây để tải file PowerPoint (.pptx) về máy", data=pptx_bytes, file_name=ten_file_pptx, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
                    
    if st.session_state.get("pdf_bytes_preview"):
        st.markdown("---")
        st.markdown("#### Bản xem trước PDF trực tiếp:")
        pdf_viewer(input=st.session_state["pdf_bytes_preview"], width=750, height=850)

# --- TAB 3: PHẢN BIỆN BỆNH ÁN ---
with tab3:
    st.markdown("### Giảng viên lâm sàng phản biện ca bệnh")
    st.caption("Giảng viên lâm sàng giàu kinh nghiệm, rà soát tính logic toàn diện của ca bệnh và đặt câu hỏi.")

    benh_su_str_pb = get_benh_su_text_for_ai()
    ca_benh_summary = f"""
    - Loại bệnh án: {loai_benh_an}
    - Bệnh sử: {benh_su_str_pb}
    - Khám toàn thân: {st.session_state.get('kham_toan_than')} (Ngày hậu phẫu: {st.session_state.get('ngay_hau_phau', 'Không có')})
    - Chẩn đoán sơ bộ: {st.session_state.get('chan_doan_so_bo')}
    - Cận lâm sàng: {st.session_state.get('cls_dx_xac_dinh')}
    - Tóm tắt: {st.session_state.get('tom_tat')}
    - Chẩn đoán xác định: {st.session_state.get('chan_doan_xac_dinh')}
    """

    col_btn_pb, col_mode = st.columns([1, 1.5])
    with col_mode: phong_cach = st.selectbox("Phong cách chất vấn:", ["Học thuật & Hướng dẫn", "Nghiêm khắc & Thách thức", "Thực chiến giao ban"])
    with col_btn_pb: st.write(""); btn_phan_bien = st.button("Giảng viên phản biện & Đặt câu hỏi", type="primary")

    if btn_phan_bien:
        model = get_feature_model("KEY_ATTENDING", "gemini-3.6-flash")
        if not model: st.error("⚠️ Hệ thống chưa cấu hình API Key!")
        elif not st.session_state.get("chan_doan_so_bo"): st.warning("⚠️ Vui lòng nhập tối thiểu Chẩn đoán sơ bộ ở Tab 1!")
        else:
            with st.spinner("Thầy/Cô đang đọc kỹ bệnh án..."):
                try:
                    prompt_phan_bien = f"""Bạn là Giảng viên lâm sàng. Hãy nhận xét ca bệnh ({ca_benh_summary}) theo phong cách {phong_cach}. 
                    Trả về đúng định dạng JSON: {{"nhan_xet_tong_the": "...", "danh_sach_cau_hoi": [{{"chu_de": "...", "cau_hoi": "...", "goi_y_tra_loi": "..."}}]}}"""
                    res_pb = model.generate_content(prompt_phan_bien)
                    res_raw = res_pb.text.strip()
                    if res_raw.startswith("```json"): res_raw = res_raw[7:]
                    elif res_raw.startswith("```"): res_raw = res_raw[3:]
                    if res_raw.endswith("```"): res_raw = res_raw[:-3]
                    st.session_state["data_phan_bien_json"] = json.loads(res_raw.strip())
                except Exception as e: st.error(f"Lỗi AI: {e}")

    if st.session_state.get("data_phan_bien_json"):
        data_pb = st.session_state["data_phan_bien_json"]
        st.divider()
        st.markdown("#### Nhận xét tổng thể & Điểm cần lưu ý:")
        st.info(data_pb.get("nhan_xet_tong_the", ""))
        st.markdown("#### ❓ Câu hỏi vấn đáp (Bấm vào từng câu để xem gợi ý đáp án):")
        for idx, item in enumerate(data_pb.get("danh_sach_cau_hoi", [])):
            with st.expander(f"**Câu {idx + 1} ({item.get('chu_de')}):** {item.get('cau_hoi')}", expanded=False):
                st.markdown("**Gợi ý hướng trả lời (Teaching Points):**")
                st.markdown(item.get("goi_y_tra_loi"))

# ==============================================================================
# CƠ CHẾ TỰ ĐỘNG LƯU NHÁP VÀO LOCALSTORAGE TRÌNH DUYỆT
# ==============================================================================
co_du_lieu = any(bool(str(st.session_state.get(k, "")).strip()) for k in ["ho_ten", "benh_su", "bs_truoc_mo", "ly_do_vao_vien", "kham_toan_than", "chan_doan_so_bo"])
if co_du_lieu:
    current_snapshot = {k: st.session_state.get(k, "") for k in FIELDS_TO_SAVE}
    for i in range(st.session_state.get("so_hang_cls", 3)):
        current_snapshot[f"cls_kq_{i}"] = st.session_state.get(f"cls_kq_{i}", "")
        current_snapshot[f"cls_pg_{i}"] = st.session_state.get(f"cls_pg_{i}", "")
    snapshot_json = json.dumps(current_snapshot, ensure_ascii=False)
    if st.session_state.get("last_saved_snapshot") != snapshot_json:
        local_storage.setItem(STORAGE_KEY, snapshot_json)
        st.session_state["last_saved_snapshot"] = snapshot_json